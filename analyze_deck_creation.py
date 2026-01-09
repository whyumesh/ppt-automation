"""
Detailed Analysis: How PowerPoint Deck is Created
Analyzes each slide to understand:
1. Which Excel file/sheet provides data
2. What operations/transformations are performed
3. How data flows from Excel to PPT
"""

import pandas as pd
from pptx import Presentation
import json
from pathlib import Path
import os

def analyze_slide_content(ppt_path, slide_num):
    """Extract detailed content from a specific slide."""
    prs = Presentation(ppt_path)
    if slide_num > len(prs.slides):
        return None
    
    slide = prs.slides[slide_num - 1]
    slide_info = {
        'slide_number': slide_num,
        'layout_name': slide.slide_layout.name,
        'shapes': []
    }
    
    for idx, shape in enumerate(slide.shapes):
        shape_info = {
            'index': idx,
            'name': shape.name,
            'type': type(shape).__name__
        }
        
        # Extract text content
        if hasattr(shape, 'text'):
            shape_info['text'] = shape.text
        
        # Extract table data
        try:
            if shape.has_table:
                table = shape.table
                shape_info['table'] = {
                    'rows': len(table.rows),
                    'columns': len(table.columns),
                    'data': []
                }
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    shape_info['table']['data'].append(row_data)
        except (ValueError, AttributeError):
            pass
        
        slide_info['shapes'].append(shape_info)
    
    return slide_info

def analyze_excel_sheet(excel_path, sheet_name, header_row=0):
    """Analyze Excel sheet structure and data."""
    try:
        df = pd.read_excel(excel_path, sheet_name=sheet_name, header=header_row)
        df = df.dropna(axis=1, how='all').dropna(axis=0, how='all')
        
        return {
            'sheet_name': sheet_name,
            'header_row': header_row,
            'total_rows': len(df),
            'columns': list(df.columns),
            'sample_data': df.head(10).to_dict(orient='records'),
            'data_types': {col: str(df[col].dtype) for col in df.columns}
        }
    except Exception as e:
        return {'error': str(e)}

def find_data_source_for_slide(slide_info, excel_data):
    """Match slide content with Excel data sources."""
    matches = []
    
    for shape in slide_info.get('shapes', []):
        if 'table' in shape:
            # Try to match table structure
            table = shape['table']
            table_data = table.get('data', [])
            
            if len(table_data) > 0:
                # Get header row
                header = table_data[0] if table_data else []
                
                # Search through Excel sheets
                for sheet_info in excel_data:
                    if 'error' in sheet_info:
                        continue
                    
                    # Check if columns match
                    excel_cols = sheet_info.get('columns', [])
                    if len(excel_cols) == len(header):
                        # Check for matching column names
                        matches_found = sum(1 for h in header if any(str(col).strip() in str(h).strip() or str(h).strip() in str(col).strip() for col in excel_cols))
                        if matches_found > 0:
                            matches.append({
                                'shape_index': shape['index'],
                                'shape_type': 'table',
                                'excel_sheet': sheet_info['sheet_name'],
                                'match_confidence': matches_found / len(header),
                                'ppt_columns': header,
                                'excel_columns': excel_cols
                            })
        
        elif 'text' in shape:
            # Try to match text content
            text = shape['text']
            for sheet_info in excel_data:
                if 'error' in sheet_info:
                    continue
                
                sample_data = sheet_info.get('sample_data', [])
                for row in sample_data:
                    for key, value in row.items():
                        if str(value) in str(text) or str(text) in str(value):
                            matches.append({
                                'shape_index': shape['index'],
                                'shape_type': 'text',
                                'excel_sheet': sheet_info['sheet_name'],
                                'matched_value': value,
                                'matched_column': key
                            })
                            break
    
    return matches

def main():
    """Main analysis function."""
    ppt_path = "Data/Apr 2025/AIL LT - April'25.pptx"
    excel_path = "Data/Apr 2025/AIL LT Working file.xlsx"
    
    print("=" * 80)
    print("DETAILED DECK CREATION ANALYSIS")
    print("=" * 80)
    print()
    
    # Get all sheets from Excel
    excel_file = pd.ExcelFile(excel_path)
    all_sheets = excel_file.sheet_names
    
    print(f"Excel File: {excel_path}")
    print(f"Total Sheets: {len(all_sheets)}")
    print(f"Sheets: {', '.join(all_sheets)}")
    print()
    
    # Analyze each sheet
    excel_data = []
    for sheet_name in all_sheets:
        print(f"Analyzing sheet: {sheet_name}...")
        # Try different header rows
        for header_row in [0, 1, 2, 3, 4]:
            sheet_info = analyze_excel_sheet(excel_path, sheet_name, header_row)
            if 'error' not in sheet_info and sheet_info['total_rows'] > 0:
                excel_data.append(sheet_info)
                break
    
    print(f"\nAnalyzed {len(excel_data)} sheets with data")
    print()
    
    # Analyze each slide
    prs = Presentation(ppt_path)
    total_slides = len(prs.slides)
    
    print("=" * 80)
    print(f"ANALYZING {total_slides} SLIDES")
    print("=" * 80)
    print()
    
    slide_analysis = []
    
    for slide_num in range(1, total_slides + 1):
        print(f"\n{'=' * 80}")
        print(f"SLIDE {slide_num}")
        print('=' * 80)
        
        slide_info = analyze_slide_content(ppt_path, slide_num)
        if not slide_info:
            continue
        
        print(f"Layout: {slide_info['layout_name']}")
        print(f"Shapes: {len(slide_info['shapes'])}")
        print()
        
        # Find data sources
        matches = find_data_source_for_slide(slide_info, excel_data)
        
        slide_analysis_item = {
            'slide_number': slide_num,
            'layout': slide_info['layout_name'],
            'shapes': slide_info['shapes'],
            'data_sources': matches
        }
        
        # Print slide details
        for shape in slide_info['shapes']:
            print(f"  Shape {shape['index']}: {shape['name']} ({shape['type']})")
            
            if 'table' in shape:
                table = shape['table']
                print(f"    Table: {table['rows']} rows x {table['columns']} columns")
                if table['data']:
                    print(f"    Headers: {table['data'][0]}")
                    if len(table['data']) > 1:
                        print(f"    Sample row: {table['data'][1]}")
            
            if 'text' in shape and shape['text']:
                text_preview = shape['text'][:100] + "..." if len(shape['text']) > 100 else shape['text']
                print(f"    Text: {text_preview}")
        
        # Print data source matches
        if matches:
            print(f"\n  [DATA SOURCES FOUND]")
            for match in matches:
                print(f"    - Excel Sheet: {match.get('excel_sheet', 'Unknown')}")
                if match.get('shape_type') == 'table':
                    print(f"      Columns: {match.get('excel_columns', [])}")
                    print(f"      Match confidence: {match.get('match_confidence', 0):.1%}")
                elif match.get('shape_type') == 'text':
                    print(f"      Matched value: {match.get('matched_value', 'N/A')}")
                    print(f"      Column: {match.get('matched_column', 'N/A')}")
        else:
            print(f"\n  [WARNING] No Excel data source found (may be static content)")
        
        slide_analysis.append(slide_analysis_item)
    
    # Save detailed analysis
    output_file = "analysis/deck_creation_analysis.json"
    os.makedirs("analysis", exist_ok=True)
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump({
            'ppt_file': ppt_path,
            'excel_file': excel_path,
            'excel_sheets_analyzed': excel_data,
            'slides_analysis': slide_analysis
        }, f, indent=2, ensure_ascii=False)
    
    print(f"\n\n{'=' * 80}")
    print(f"Analysis saved to: {output_file}")
    print('=' * 80)

if __name__ == "__main__":
    main()
