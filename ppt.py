"""
Automated PowerPoint Generator for Employee Performance Data
Author: Your Name
Version: 1.0
Date: 2025-10-08
"""

import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import datetime as dt
import os

# ====== CONFIGURATION ======
INPUT_FILE = "employee_performance_data.xlsx"  # your Excel file name
OUTPUT_FILE = f"Employee_Performance_Report_{dt.datetime.now().strftime('%Y%m%d_%H%M')}.pptx"

# Colors and theme
TITLE_COLOR = RGBColor(0, 51, 102)
TEXT_COLOR = RGBColor(60, 60, 60)
CHART_COLORS = ['#1f77b4', '#ff7f0e']
# ===========================


def create_styled_ppt(df):
    prs = Presentation()

    # ---------- Summary Slide ----------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "Employee Performance Report"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.color.rgb = TITLE_COLOR

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
    sf = sub_box.text_frame
    sf.text = f"Generated on: {dt.datetime.now().strftime('%d %B %Y, %I:%M %p')}"
    sf.paragraphs[0].font.size = Pt(20)
    sf.paragraphs[0].alignment = PP_ALIGN.CENTER
    sf.paragraphs[0].font.color.rgb = TEXT_COLOR

    # Department-wise summary chart
    dept_summary = df.groupby("Department")[["Sales", "Rating"]].mean().reset_index()
    fig, ax = plt.subplots(figsize=(5, 3))
    ax.bar(dept_summary["Department"], dept_summary["Sales"], color=CHART_COLORS[0])
    ax.set_title("Average Sales by Department", fontsize=12)
    plt.tight_layout()

    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=150)
    plt.close(fig)
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, Inches(2), Inches(3), Inches(6), Inches(3))

    # ---------- Individual Employee Slides ----------
    max_sales = df["Sales"].max()
    max_rating = df["Rating"].max()

    for idx, row in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Employee name (title)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = str(row["Name"])
        title_frame.paragraphs[0].font.size = Pt(34)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

        # Department & Month (subheading)
        sub_text = f"Department: {row['Department']}   |   Month: {row['Month']}"
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
        sub_frame = sub_box.text_frame
        sub_frame.text = sub_text
        sub_frame.paragraphs[0].font.size = Pt(18)
        sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        sub_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

        # Sales chart
        fig, ax = plt.subplots(figsize=(4, 3))
        ax.bar(["Sales", "Max Sales"], [row["Sales"], max_sales], color=CHART_COLORS)
        ax.set_title("Sales Performance", fontsize=12)
        for i, v in enumerate([row["Sales"], max_sales]):
            ax.text(i, v + max_sales * 0.02, f"{v:.0f}", ha='center', fontweight='bold')
        plt.tight_layout()

        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=150)
        plt.close(fig)
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(2), Inches(4), Inches(3))

        # Rating chart
        fig, ax = plt.subplots(figsize=(4, 3))
        ax.bar(["Rating"], [row["Rating"]], color=CHART_COLORS[1])
        ax.set_ylim(0, max_rating + 1)
        ax.set_title("Rating", fontsize=12)
        ax.text(0, row["Rating"] + 0.1, f"{row['Rating']:.1f}", ha='center', fontweight='bold')
        plt.tight_layout()

        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=150)
        plt.close(fig)
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(5), Inches(2), Inches(4), Inches(3))

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.5))
        ff = footer_box.text_frame
        ff.text = f"Generated on {dt.datetime.now().strftime('%d %b %Y')}"
        ff.paragraphs[0].font.size = Pt(12)
        ff.paragraphs[0].alignment = PP_ALIGN.RIGHT
        ff.paragraphs[0].font.color.rgb = TEXT_COLOR

    # Save PowerPoint
    prs.save(OUTPUT_FILE)
    print(f"✅ PowerPoint created: {OUTPUT_FILE}")


# ---------- MAIN EXECUTION ----------
if __name__ == "__main__":
    if not os.path.exists(INPUT_FILE):
        raise FileNotFoundError(f"❌ Excel file not found: {INPUT_FILE}")

    df = pd.read_excel(INPUT_FILE)
    print(f"📊 Loaded {len(df)} records with columns: {list(df.columns)}")
    create_styled_ppt(df)
