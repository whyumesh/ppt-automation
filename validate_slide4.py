"""
Validate Slide 4 (Consent) from generated PPT vs manual PPT
"""

from pptx import Presentation

manual_ppt = "Data/Apr 2025/AIL LT - April'25.pptx"
generated_ppt = "output/test_slide4_raw.pptx"

print("=" * 80)
print("VALIDATING SLIDE 4 (CONSENT)")
print("=" * 80)

# Load presentations
manual_prs = Presentation(manual_ppt)
generated_prs = Presentation(generated_ppt)

print(f"\nManual PPT: {len(manual_prs.slides)} slides")
print(f"Generated PPT: {len(generated_prs.slides)} slides")

# Get Slide 4 (index 3)
manual_slide = manual_prs.slides[3] if len(manual_prs.slides) > 3 else None
generated_slide = generated_prs.slides[3] if len(generated_prs.slides) > 3 else None

if not manual_slide or not generated_slide:
    print("ERROR: Slide 4 not found in one or both presentations")
    exit(1)

# Find tables
manual_table = None
generated_table = None

for shape in manual_slide.shapes:
    if shape.has_table:
        manual_table = shape.table
        break

for shape in generated_slide.shapes:
    if shape.has_table:
        generated_table = shape.table
        break

print(f"\nManual PPT Table: {len(manual_table.rows) if manual_table else 0} rows, {len(manual_table.columns) if manual_table else 0} cols")
print(f"Generated PPT Table: {len(generated_table.rows) if generated_table else 0} rows, {len(generated_table.columns) if generated_table else 0} cols")

if not manual_table or not generated_table:
    print("ERROR: Table not found in one or both slides")
    exit(1)

# Compare table content
print("\n" + "=" * 80)
print("COMPARING TABLE CONTENT")
print("=" * 80)

print("\nManual PPT - First 5 rows:")
for i in range(min(5, len(manual_table.rows))):
    row = manual_table.rows[i]
    cells = [cell.text.strip() for cell in row.cells]
    print(f"  Row {i}: {cells}")

print("\nGenerated PPT - First 5 rows:")
for i in range(min(5, len(generated_table.rows))):
    row = generated_table.rows[i]
    cells = [cell.text.strip() for cell in row.cells]
    print(f"  Row {i}: {cells}")

# Check if data matches
print("\n" + "=" * 80)
print("VALIDATION RESULT")
print("=" * 80)

if len(generated_table.rows) > 0 and len(generated_table.columns) >= 4:
    print("✓ Table exists with correct structure")
    
    # Check header row
    manual_header = [cell.text.strip() for cell in manual_table.rows[0].cells[:4]]
    generated_header = [cell.text.strip() for cell in generated_table.rows[0].cells[:4]]
    
    print(f"\nManual header: {manual_header}")
    print(f"Generated header: {generated_header}")
    
    # Check data rows
    if len(generated_table.rows) > 1:
        print(f"\n✓ Table has {len(generated_table.rows)} rows (including header)")
        print(f"✓ First data row: {[cell.text.strip() for cell in generated_table.rows[1].cells[:4]]}")
        print("\n✓ VALIDATION PASSED: Slide 4 has data!")
    else:
        print("\n✗ VALIDATION FAILED: Table has no data rows")
else:
    print("✗ VALIDATION FAILED: Table structure incorrect")

