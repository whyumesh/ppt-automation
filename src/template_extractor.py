"""
PPT Template Extractor
Analyzes existing PowerPoint files to extract slide layouts, placeholders, and formatting rules.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import json
import os
from typing import Dict, List, Any, Optional, Tuple


def extract_rgb_from_color(rgb_color) -> Optional[Dict[str, int]]:
    """
    Safely extract RGB values from an RGBColor object.
    
    Args:
        rgb_color: RGBColor object (tuple-like) or integer
    
    Returns:
        Dictionary with r, g, b values or None
    """
    if rgb_color is None:
        return None
    
    try:
        # RGBColor is tuple-like: [0]=red, [1]=green, [2]=blue
        if hasattr(rgb_color, '__getitem__') and hasattr(rgb_color, '__len__'):
            if len(rgb_color) >= 3:
                return {
                    "r": rgb_color[0],
                    "g": rgb_color[1],
                    "b": rgb_color[2]
                }
        # If it's an integer, extract RGB components
        elif isinstance(rgb_color, int):
            return {
                "r": (rgb_color >> 16) & 0xFF,
                "g": (rgb_color >> 8) & 0xFF,
                "b": rgb_color & 0xFF
            }
    except (TypeError, AttributeError, ValueError, IndexError):
        pass
    
    return None


class TemplateExtractor:
    """Extracts template information from existing PowerPoint files."""
    
    def __init__(self, ppt_path: str):
        """
        Initialize the template extractor.
        
        Args:
            ppt_path: Path to the PowerPoint file to analyze
        """
        if not os.path.exists(ppt_path):
            raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")
        
        self.ppt_path = ppt_path
        self.presentation = Presentation(ppt_path)
        self.template_info = {
            "file_path": ppt_path,
            "slide_count": len(self.presentation.slides),
            "slides": []
        }
    
    def extract_all(self) -> Dict[str, Any]:
        """
        Extract all template information from the presentation.
        
        Returns:
            Dictionary containing template structure and formatting information
        """
        for idx, slide in enumerate(self.presentation.slides):
            slide_info = self.extract_slide_info(slide, idx)
            self.template_info["slides"].append(slide_info)
        
        return self.template_info
    
    def extract_slide_info(self, slide, slide_index: int) -> Dict[str, Any]:
        """
        Extract information from a single slide.
        
        Args:
            slide: PowerPoint slide object
            slide_index: Index of the slide (0-based)
        
        Returns:
            Dictionary containing slide information
        """
        slide_info = {
            "slide_number": slide_index + 1,
            "slide_id": slide.slide_id,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "shapes": []
        }
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = self.extract_shape_info(shape, shape_idx)
            if shape_info:
                slide_info["shapes"].append(shape_info)
        
        return slide_info
    
    def extract_shape_info(self, shape, shape_index: int) -> Optional[Dict[str, Any]]:
        """
        Extract information from a shape (text box, table, chart, etc.).
        
        Args:
            shape: PowerPoint shape object
            shape_index: Index of the shape on the slide
        
        Returns:
            Dictionary containing shape information, or None if shape is not extractable
        """
        shape_info = {
            "index": shape_index,
            "shape_id": shape.shape_id,
            "name": shape.name,
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "shape_type": str(shape.shape_type)
        }
        
        # Extract text box information
        if shape.has_text_frame:
            text_frame = shape.text_frame
            shape_info["type"] = "text_box"
            shape_info["text_content"] = []
            shape_info["formatting"] = {}
            
            for paragraph in text_frame.paragraphs:
                para_info = {
                    "text": paragraph.text,
                    "level": paragraph.level,
                    "alignment": str(paragraph.alignment) if paragraph.alignment else None,
                    "runs": []
                }
                
                for run in paragraph.runs:
                    run_info = {
                        "text": run.text,
                        "font_name": run.font.name if run.font.name else None,
                        "font_size": run.font.size.pt if run.font.size else None,
                        "bold": run.font.bold,
                        "italic": run.font.italic,
                        "underline": run.font.underline,
                    }
                    
                    # Extract font color
                    if run.font.color:
                        try:
                            # Check if it's an RGB color type
                            if run.font.color.type == 1:  # RGB color
                                color = run.font.color
                                if hasattr(color, 'rgb') and color.rgb:
                                    rgb_dict = extract_rgb_from_color(color.rgb)
                                    if rgb_dict:
                                        run_info["font_color"] = rgb_dict
                        except AttributeError:
                            # Color is not RGB type (e.g., scheme color) - skip
                            pass
                    
                    para_info["runs"].append(run_info)
                
                shape_info["text_content"].append(para_info)
            
            # Extract text frame formatting
            if text_frame.auto_size:
                shape_info["formatting"]["auto_size"] = str(text_frame.auto_size)
            shape_info["formatting"]["margin_left"] = text_frame.margin_left
            shape_info["formatting"]["margin_right"] = text_frame.margin_right
            shape_info["formatting"]["margin_top"] = text_frame.margin_top
            shape_info["formatting"]["margin_bottom"] = text_frame.margin_bottom
        
        # Extract table information
        elif shape.has_table:
            table = shape.table
            shape_info["type"] = "table"
            shape_info["table_info"] = {
                "rows": len(table.rows),
                "columns": len(table.columns),
                "cells": []
            }
            
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    cell_info = {
                        "row": row_idx,
                        "column": col_idx,
                        "text": cell.text,
                        "fill_color": None
                    }
                    
                    # Extract cell fill color
                    if cell.fill:
                        try:
                            # Try to access fore_color (may raise TypeError for NoFill types)
                            fore_color = cell.fill.fore_color
                            if fore_color:
                                # Try to get RGB color (may not exist for scheme colors)
                                try:
                                    rgb_color = fore_color.rgb
                                    if rgb_color:
                                        rgb_dict = extract_rgb_from_color(rgb_color)
                                        if rgb_dict:
                                            cell_info["fill_color"] = rgb_dict
                                except AttributeError:
                                    # Color is not RGB type (e.g., scheme color) - skip
                                    pass
                        except (AttributeError, TypeError):
                            # Fill type doesn't support foreground color (e.g., NoFill) - skip
                            pass
                    
                    shape_info["table_info"]["cells"].append(cell_info)
        
        # Extract chart information
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            shape_info["type"] = "chart"
            shape_info["chart_type"] = str(shape.chart.chart_type) if hasattr(shape, 'chart') else None
        
        # Extract picture information
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_info["type"] = "picture"
            shape_info["image_format"] = shape.image.ext if hasattr(shape, 'image') else None
        
        else:
            shape_info["type"] = "other"
        
        return shape_info
    
    def save_template_info(self, output_path: str):
        """
        Save extracted template information to a JSON file.
        
        Args:
            output_path: Path to save the JSON file
        """
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(self.template_info, f, indent=2, ensure_ascii=False)
    
    def create_template_copy(self, output_path: str):
        """
        Create a clean template copy of the presentation (without content).
        
        Args:
            output_path: Path to save the template file
        """
        template_prs = Presentation()
        
        # Copy slide layouts (this is a simplified version)
        # In practice, we'd need to copy layouts more carefully
        for slide in self.presentation.slides:
            new_slide = template_prs.slides.add_slide(slide.slide_layout)
            
            # Copy shape structure but clear text content
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Keep structure but clear text
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.text = ""
        
        template_prs.save(output_path)


def extract_template_from_file(ppt_path: str, output_json: Optional[str] = None, 
                              output_template: Optional[str] = None) -> Dict[str, Any]:
    """
    Convenience function to extract template information from a PowerPoint file.
    
    Args:
        ppt_path: Path to the PowerPoint file
        output_json: Optional path to save JSON output
        output_template: Optional path to save template copy
    
    Returns:
        Dictionary containing template information
    """
    extractor = TemplateExtractor(ppt_path)
    template_info = extractor.extract_all()
    
    if output_json:
        extractor.save_template_info(output_json)
    
    if output_template:
        extractor.create_template_copy(output_template)
    
    return template_info


if __name__ == "__main__":
    # Example usage
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python template_extractor.py <ppt_file> [output_json] [output_template]")
        sys.exit(1)
    
    ppt_file = sys.argv[1]
    output_json = sys.argv[2] if len(sys.argv) > 2 else None
    output_template = sys.argv[3] if len(sys.argv) > 3 else None
    
    info = extract_template_from_file(ppt_file, output_json, output_template)
    print(f"Extracted template information from {ppt_file}")
    print(f"Found {info['slide_count']} slides")

