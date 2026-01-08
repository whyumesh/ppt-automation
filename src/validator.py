"""
PPT Validator
Compares generated PPTs with manual versions and creates validation reports.
"""

import os
import sys
from pptx import Presentation
from typing import Dict, List, Any, Optional, Tuple
import json
from datetime import datetime

# Handle imports for both module and script usage
try:
    from template_extractor import TemplateExtractor
except ImportError:
    try:
        from src.template_extractor import TemplateExtractor
    except ImportError:
        # If running as script, add src to path
        sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))
        from src.template_extractor import TemplateExtractor


class PPTValidator:
    """Validates generated PowerPoint decks against manual versions."""
    
    def __init__(self, manual_ppt_path: str, generated_ppt_path: str):
        """
        Initialize the validator.
        
        Args:
            manual_ppt_path: Path to manually created PowerPoint file
            generated_ppt_path: Path to generated PowerPoint file
        """
        if not os.path.exists(manual_ppt_path):
            raise FileNotFoundError(f"Manual PPT file not found: {manual_ppt_path}")
        if not os.path.exists(generated_ppt_path):
            raise FileNotFoundError(f"Generated PPT file not found: {generated_ppt_path}")
        
        self.manual_ppt_path = manual_ppt_path
        self.generated_ppt_path = generated_ppt_path
        
        self.manual_prs = Presentation(manual_ppt_path)
        self.generated_prs = Presentation(generated_ppt_path)
        
        self.validation_results = {
            "manual_file": manual_ppt_path,
            "generated_file": generated_ppt_path,
            "validation_date": datetime.now().isoformat(),
            "slide_count_match": False,
            "slides": [],
            "summary": {
                "total_slides": 0,
                "matching_slides": 0,
                "mismatches": 0,
                "accuracy": 0.0
            }
        }
    
    def validate_all(self) -> Dict[str, Any]:
        """
        Validate all slides in both presentations.
        
        Returns:
            Dictionary containing validation results
        """
        # Check slide count
        manual_count = len(self.manual_prs.slides)
        generated_count = len(self.generated_prs.slides)
        
        self.validation_results["slide_count_match"] = manual_count == generated_count
        self.validation_results["summary"]["total_slides"] = manual_count
        
        # Validate each slide
        max_slides = min(manual_count, generated_count)
        
        for slide_idx in range(max_slides):
            slide_result = self.validate_slide(slide_idx)
            self.validation_results["slides"].append(slide_result)
            
            if slide_result["match"]:
                self.validation_results["summary"]["matching_slides"] += 1
            else:
                self.validation_results["summary"]["mismatches"] += 1
        
        # Calculate accuracy
        if max_slides > 0:
            self.validation_results["summary"]["accuracy"] = (
                self.validation_results["summary"]["matching_slides"] / max_slides * 100
            )
        
        return self.validation_results
    
    def validate_slide(self, slide_index: int) -> Dict[str, Any]:
        """
        Validate a single slide.
        
        Args:
            slide_index: Index of the slide to validate (0-based)
        
        Returns:
            Dictionary containing slide validation results
        """
        manual_slide = self.manual_prs.slides[slide_index]
        generated_slide = self.generated_prs.slides[slide_index]
        
        slide_result = {
            "slide_number": slide_index + 1,
            "match": True,
            "shape_count_match": False,
            "shapes": [],
            "errors": []
        }
        
        # Check shape count
        manual_shape_count = len(manual_slide.shapes)
        generated_shape_count = len(generated_slide.shapes)
        slide_result["shape_count_match"] = manual_shape_count == generated_shape_count
        
        # Validate shapes
        max_shapes = min(manual_shape_count, generated_shape_count)
        
        for shape_idx in range(max_shapes):
            shape_result = self.validate_shape(
                manual_slide.shapes[shape_idx],
                generated_slide.shapes[shape_idx],
                shape_idx
            )
            slide_result["shapes"].append(shape_result)
            
            if not shape_result["match"]:
                slide_result["match"] = False
        
        # If shape counts don't match, add error
        if not slide_result["shape_count_match"]:
            slide_result["errors"].append(
                f"Shape count mismatch: manual={manual_shape_count}, generated={generated_shape_count}"
            )
            slide_result["match"] = False
        
        return slide_result
    
    def validate_shape(self, manual_shape, generated_shape, shape_index: int) -> Dict[str, Any]:
        """
        Validate a single shape.
        
        Args:
            manual_shape: Manual shape object
            generated_shape: Generated shape object
            shape_index: Index of the shape
        
        Returns:
            Dictionary containing shape validation results
        """
        shape_result = {
            "shape_index": shape_index,
            "match": True,
            "type_match": False,
            "text_match": False,
            "errors": []
        }
        
        # Check shape type
        manual_type = str(manual_shape.shape_type)
        generated_type = str(generated_shape.shape_type)
        shape_result["type_match"] = manual_type == generated_type
        
        if not shape_result["type_match"]:
            shape_result["errors"].append(
                f"Type mismatch: manual={manual_type}, generated={generated_type}"
            )
            shape_result["match"] = False
        
        # Check text content
        if manual_shape.has_text_frame and generated_shape.has_text_frame:
            manual_text = manual_shape.text_frame.text
            generated_text = generated_shape.text_frame.text
            
            # Normalize text for comparison
            manual_text_normalized = self._normalize_text(manual_text)
            generated_text_normalized = self._normalize_text(generated_text)
            
            shape_result["text_match"] = manual_text_normalized == generated_text_normalized
            
            if not shape_result["text_match"]:
                shape_result["errors"].append(
                    f"Text mismatch:\n  Manual: {manual_text[:100]}...\n  Generated: {generated_text[:100]}..."
                )
                shape_result["match"] = False
        
        # Check table content
        elif manual_shape.has_table and generated_shape.has_table:
            table_result = self.validate_table(manual_shape.table, generated_shape.table)
            shape_result["table_match"] = table_result["match"]
            shape_result["table_errors"] = table_result.get("errors", [])
            
            if not table_result["match"]:
                shape_result["match"] = False
        
        return shape_result
    
    def validate_table(self, manual_table, generated_table) -> Dict[str, Any]:
        """
        Validate a table.
        
        Args:
            manual_table: Manual table object
            generated_table: Generated table object
        
        Returns:
            Dictionary containing table validation results
        """
        table_result = {
            "match": True,
            "row_count_match": False,
            "column_count_match": False,
            "cell_matches": [],
            "errors": []
        }
        
        manual_rows = len(manual_table.rows)
        generated_rows = len(generated_table.rows)
        manual_cols = len(manual_table.columns)
        generated_cols = len(generated_table.columns)
        
        table_result["row_count_match"] = manual_rows == generated_rows
        table_result["column_count_match"] = manual_cols == generated_cols
        
        if not table_result["row_count_match"]:
            table_result["errors"].append(
                f"Row count mismatch: manual={manual_rows}, generated={generated_rows}"
            )
            table_result["match"] = False
        
        if not table_result["column_count_match"]:
            table_result["errors"].append(
                f"Column count mismatch: manual={manual_cols}, generated={generated_cols}"
            )
            table_result["match"] = False
        
        # Validate cell content
        max_rows = min(manual_rows, generated_rows)
        max_cols = min(manual_cols, generated_cols)
        
        for row_idx in range(max_rows):
            for col_idx in range(max_cols):
                manual_cell_text = manual_table.cell(row_idx, col_idx).text
                generated_cell_text = generated_table.cell(row_idx, col_idx).text
                
                manual_normalized = self._normalize_text(manual_cell_text)
                generated_normalized = self._normalize_text(generated_cell_text)
                
                cell_match = manual_normalized == generated_normalized
                table_result["cell_matches"].append({
                    "row": row_idx,
                    "column": col_idx,
                    "match": cell_match,
                    "manual": manual_cell_text,
                    "generated": generated_cell_text
                })
                
                if not cell_match:
                    table_result["match"] = False
        
        return table_result
    
    def _normalize_text(self, text: str) -> str:
        """
        Normalize text for comparison.
        
        Args:
            text: Text to normalize
        
        Returns:
            Normalized text
        """
        # Remove extra whitespace
        text = " ".join(text.split())
        # Convert to lowercase for comparison
        text = text.lower()
        # Remove special characters that might differ
        text = text.replace("\n", " ").replace("\r", " ")
        return text.strip()
    
    def save_report(self, output_path: str):
        """
        Save validation report to a JSON file.
        
        Args:
            output_path: Path to save the report
        """
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(self.validation_results, f, indent=2, ensure_ascii=False)
    
    def print_summary(self):
        """Print a summary of validation results."""
        summary = self.validation_results["summary"]
        
        print("\n" + "="*60)
        print("VALIDATION SUMMARY")
        print("="*60)
        print(f"Total Slides: {summary['total_slides']}")
        print(f"Matching Slides: {summary['matching_slides']}")
        print(f"Mismatches: {summary['mismatches']}")
        print(f"Accuracy: {summary['accuracy']:.2f}%")
        print("="*60)
        
        # Print slide-level errors
        for slide_result in self.validation_results["slides"]:
            if not slide_result["match"]:
                print(f"\nSlide {slide_result['slide_number']} - MISMATCH:")
                for error in slide_result.get("errors", []):
                    print(f"  - {error}")
                
                # Print shape-level errors
                for shape_result in slide_result.get("shapes", []):
                    if not shape_result["match"]:
                        print(f"  Shape {shape_result['shape_index']}:")
                        for error in shape_result.get("errors", []):
                            print(f"    - {error}")


def validate_ppt(manual_ppt_path: str, generated_ppt_path: str, 
                output_report: Optional[str] = None) -> Dict[str, Any]:
    """
    Convenience function to validate a generated PPT against a manual version.
    
    Args:
        manual_ppt_path: Path to manually created PowerPoint file
        generated_ppt_path: Path to generated PowerPoint file
        output_report: Optional path to save validation report
    
    Returns:
        Dictionary containing validation results
    """
    validator = PPTValidator(manual_ppt_path, generated_ppt_path)
    results = validator.validate_all()
    
    validator.print_summary()
    
    if output_report:
        validator.save_report(output_report)
        print(f"\nValidation report saved to: {output_report}")
    
    return results


if __name__ == "__main__":
    # Example usage
    import sys
    
    if len(sys.argv) < 3:
        print("Usage: python validator.py <manual_ppt> <generated_ppt> [output_report]")
        sys.exit(1)
    
    manual_ppt = sys.argv[1]
    generated_ppt = sys.argv[2]
    output_report = sys.argv[3] if len(sys.argv) > 3 else None
    
    validate_ppt(manual_ppt, generated_ppt, output_report)

