"""
PPT Formatter
Applies formatting (fonts, colors, alignment) to PowerPoint elements.
"""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Pt
from typing import Dict, Any, Optional, Tuple


class PPTFormatter:
    """Applies formatting to PowerPoint elements."""
    
    def __init__(self, formatting_config: Optional[Dict] = None):
        """
        Initialize the formatter.
        
        Args:
            formatting_config: Optional formatting configuration dictionary
        """
        self.formatting_config = formatting_config or {}
        self.default_font_size = self.formatting_config.get("fonts", {}).get("default_size", 12)
        self.default_font_name = self.formatting_config.get("fonts", {}).get("default_name", "Calibri")
    
    def format_text_box(self, text_frame, formatting: Dict[str, Any]):
        """
        Format a text frame.
        
        Args:
            text_frame: PowerPoint text frame object
            formatting: Dictionary containing formatting options
        """
        # Set margins
        if "margin_left" in formatting:
            text_frame.margin_left = formatting["margin_left"]
        if "margin_right" in formatting:
            text_frame.margin_right = formatting["margin_right"]
        if "margin_top" in formatting:
            text_frame.margin_top = formatting["margin_top"]
        if "margin_bottom" in formatting:
            text_frame.margin_bottom = formatting["margin_bottom"]
        
        # Format paragraphs
        for paragraph in text_frame.paragraphs:
            self.format_paragraph(paragraph, formatting)
    
    def format_paragraph(self, paragraph, formatting: Dict[str, Any]):
        """
        Format a paragraph.
        
        Args:
            paragraph: PowerPoint paragraph object
            formatting: Dictionary containing formatting options
        """
        # Set alignment
        if "alignment" in formatting:
            alignment_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY
            }
            alignment = alignment_map.get(formatting["alignment"].lower(), PP_ALIGN.LEFT)
            paragraph.alignment = alignment
        
        # Format runs
        for run in paragraph.runs:
            self.format_text_run(run, formatting)
    
    def format_text_run(self, run, formatting: Dict[str, Any]):
        """
        Format a text run.
        
        Args:
            run: PowerPoint text run object
            formatting: Dictionary containing formatting options
        """
        font = run.font
        
        # Set font name
        if "font_name" in formatting:
            font.name = formatting["font_name"]
        elif self.default_font_name:
            font.name = self.default_font_name
        
        # Set font size
        if "font_size" in formatting:
            font.size = Pt(formatting["font_size"])
        elif self.default_font_size:
            font.size = Pt(self.default_font_size)
        
        # Set font style
        if "bold" in formatting:
            font.bold = formatting["bold"]
        if "italic" in formatting:
            font.italic = formatting["italic"]
        if "underline" in formatting:
            font.underline = formatting["underline"]
        
        # Set font color
        if "font_color" in formatting:
            color = formatting["font_color"]
            if isinstance(color, str):
                # Hex color string
                if color.startswith("#"):
                    color = color[1:]
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                font.color.rgb = RGBColor(r, g, b)
            elif isinstance(color, dict):
                # RGB dictionary
                r = color.get("r", 0)
                g = color.get("g", 0)
                b = color.get("b", 0)
                font.color.rgb = RGBColor(r, g, b)
    
    def format_table_cell(self, cell, formatting: Dict[str, Any]):
        """
        Format a table cell.
        
        Args:
            cell: PowerPoint table cell object
            formatting: Dictionary containing formatting options
        """
        # Set cell fill color
        if "fill_color" in formatting:
            color = formatting["fill_color"]
            if isinstance(color, str):
                if color.startswith("#"):
                    color = color[1:]
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            elif isinstance(color, dict):
                r = color.get("r", 0)
                g = color.get("g", 0)
                b = color.get("b", 0)
            else:
                return
            
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(r, g, b)
        
        # Format text in cell
        if cell.text_frame:
            self.format_text_box(cell.text_frame, formatting)
    
    def format_table(self, table, formatting: Dict[str, Any]):
        """
        Format a table.
        
        Args:
            table: PowerPoint table object
            formatting: Dictionary containing formatting options
        """
        # Format header row if specified
        if "header_formatting" in formatting and len(table.rows) > 0:
            header_row = table.rows[0]
            for cell in header_row.cells:
                self.format_table_cell(cell, formatting["header_formatting"])
        
        # Format data rows if specified
        if "data_formatting" in formatting:
            for row in table.rows[1:]:
                for cell in row.cells:
                    self.format_table_cell(cell, formatting["data_formatting"])
        
        # Apply general formatting to all cells
        if "cell_formatting" in formatting:
            for row in table.rows:
                for cell in row.cells:
                    self.format_table_cell(cell, formatting["cell_formatting"])
    
    def apply_conditional_formatting(self, element, value: float, 
                                    threshold: float = 0.0,
                                    positive_formatting: Optional[Dict] = None,
                                    negative_formatting: Optional[Dict] = None):
        """
        Apply conditional formatting based on value and threshold.
        
        Args:
            element: PowerPoint element to format (text frame, cell, etc.)
            value: Value to evaluate
            threshold: Threshold for conditional formatting
            positive_formatting: Formatting to apply if value >= threshold
            negative_formatting: Formatting to apply if value < threshold
        """
        if value >= threshold:
            formatting = positive_formatting or {}
        else:
            formatting = negative_formatting or {}
        
        if isinstance(element, type) and hasattr(element, 'text_frame'):
            self.format_text_box(element.text_frame, formatting)
        elif hasattr(element, 'text_frame'):
            self.format_text_box(element.text_frame, formatting)


def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """
    Convert hex color string to RGB tuple.
    
    Args:
        hex_color: Hex color string (with or without #)
    
    Returns:
        RGB tuple (r, g, b)
    """
    if hex_color.startswith("#"):
        hex_color = hex_color[1:]
    
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    
    return (r, g, b)


def rgb_to_hex(rgb: Tuple[int, int, int]) -> str:
    """
    Convert RGB tuple to hex color string.
    
    Args:
        rgb: RGB tuple (r, g, b)
    
    Returns:
        Hex color string
    """
    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"

