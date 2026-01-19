"""
PPT Builder
Slide building utilities for creating PowerPoint slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from typing import Dict, List, Any, Optional
import pandas as pd
try:
    from .ppt_formatter import PPTFormatter
except ImportError:
    from ppt_formatter import PPTFormatter


class PPTBuilder:
    """Builds PowerPoint slides from data."""
    
    def __init__(self, presentation: Presentation, formatter: Optional[PPTFormatter] = None):
        """
        Initialize the PPT builder.
        
        Args:
            presentation: PowerPoint presentation object
            formatter: Optional PPTFormatter instance
        """
        self.presentation = presentation
        self.formatter = formatter or PPTFormatter()
    
    def add_slide(self, layout=None) -> Any:
        """
        Add a new slide to the presentation.
        
        Args:
            layout: Optional slide layout to use
        
        Returns:
            New slide object
        """
        if layout:
            slide = self.presentation.slides.add_slide(layout)
        else:
            slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        
        return slide
    
    def add_text_box(self, slide, text: str, left: float, top: float,
                    width: float, height: float, formatting: Optional[Dict] = None) -> Any:
        """
        Add a text box to a slide.
        
        Args:
            slide: Slide object
            text: Text content
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            formatting: Optional formatting dictionary
        
        Returns:
            Shape object
        """
        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(width)
        height_inches = Inches(height)
        
        text_box = slide.shapes.add_textbox(left_inches, top_inches, width_inches, height_inches)
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        text_frame.margin_bottom = Inches(0.05)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        
        # Apply default formatting if none specified
        if formatting is None:
            formatting = {
                "font_size": 18,
                "font_name": "Calibri",
                "alignment": "left"
            }
        
        if formatting:
            self.formatter.format_text_box(text_frame, formatting)
        
        return text_box
    
    def add_table(self, slide, data: pd.DataFrame, left: float, top: float,
                 width: float, height: float, formatting: Optional[Dict] = None) -> Any:
        """
        Add a table to a slide with fully adaptive sizing and row type detection.
        
        Args:
            slide: Slide object
            data: DataFrame containing table data
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            formatting: Optional formatting dictionary
        
        Returns:
            Shape object
        """
        # Handle empty DataFrame - create table with headers and "No data" message
        if data is None:
            # Create empty DataFrame with default structure
            data = pd.DataFrame({"Column": ["No data available"]})
        
        if len(data) == 0:
            # Create a single row with "No data available" message
            if len(data.columns) == 0:
                # No columns at all - create a default column
                data = pd.DataFrame({"Data": ["No data available"]})
            else:
                # Has columns but no rows - add one row with message
                no_data_row = {col: "No data available" for col in data.columns}
                data = pd.DataFrame([no_data_row])
        
        rows = len(data) + 1  # +1 for header
        cols = len(data.columns)
        
        # STANDARD SLIDE DIMENSIONS (PowerPoint)
        SLIDE_HEIGHT = 7.5  # inches
        SLIDE_WIDTH = 10.0  # inches
        
        # Calculate available space - ensure we don't exceed slide bounds
        # The height parameter is a MAXIMUM constraint - we'll calculate optimal height below
        available_width = min(width, SLIDE_WIDTH - left - 0.5)  # Leave right margin
        max_available_height = min(height, SLIDE_HEIGHT - top - 0.5)  # Maximum height constraint
        
        # ADAPTIVE SYSTEM: Calculate optimal sizes based on data
        num_data_rows = len(data)
        
        # Detect row types BEFORE creating table
        row_types = []  # 'regular', 'subtotal', 'total'
        for idx, row in data.iterrows():
            row_type = 'regular'
            
            # Check first column for patterns indicating subtotal/total rows
            first_col_value = str(row.iloc[0]).strip().upper() if len(row) > 0 else ""
            second_col_value = ""
            if len(row) > 1:
                second_val = row.iloc[1]
                if pd.isna(second_val):
                    second_col_value = ""
                else:
                    second_col_value = str(second_val).strip()
            
            # Total row detection: Contains "TOTAL" and is likely last row or grand total
            if 'TOTAL' in first_col_value:
                # Check if it's a grand total (contains company name like "AIL" or starts with "TOTAL" or "GRAND")
                if first_col_value.startswith(('AIL', 'TOTAL', 'GRAND')):
                    row_type = 'total'
                elif idx == len(data) - 1:
                    # Last row with "TOTAL" is likely grand total
                    row_type = 'total'
                else:
                    # Could be a subtotal with "Total" in name
                    row_type = 'subtotal'
            # Subtotal row detection: First column has name, second column is empty/NaN
            elif (pd.isna(row.iloc[1]) if len(row) > 1 else True) or second_col_value == "":
                # If first column has a name (not empty) and second is empty, likely subtotal
                if first_col_value and first_col_value != "" and first_col_value != "NAN":
                    # Check if it's not the first row and has data in other columns
                    has_data = False
                    for i in range(2, min(len(row), len(data.columns))):
                        if not pd.isna(row.iloc[i]) and str(row.iloc[i]).strip() != "":
                            has_data = True
                            break
                    if has_data:
                        row_type = 'subtotal'
            
            row_types.append(row_type)
        
        # Count row types for adaptive sizing
        num_subtotal_rows = sum(1 for rt in row_types if rt == 'subtotal')
        num_total_rows = sum(1 for rt in row_types if rt == 'total')
        num_regular_rows = num_data_rows - num_subtotal_rows - num_total_rows
        
        print(f"DEBUG: Row type detection - Regular: {num_regular_rows}, Subtotal: {num_subtotal_rows}, Total: {num_total_rows}")
        
        # ADAPTIVE ROW HEIGHTS - Calculate BEFORE creating table
        if num_data_rows <= 5:
            base_row_height = 0.50
            base_header_height = 0.55
        elif num_data_rows <= 10:
            base_row_height = 0.42
            base_header_height = 0.47
        elif num_data_rows <= 15:
            base_row_height = 0.35
            base_header_height = 0.40
        elif num_data_rows <= 20:
            base_row_height = 0.30
            base_header_height = 0.35
        elif num_data_rows <= 25:
            base_row_height = 0.26
            base_header_height = 0.30
        elif num_data_rows <= 30:
            base_row_height = 0.23
            base_header_height = 0.27
        else:
            base_row_height = 0.20
            base_header_height = 0.24
        
        # Calculate total height needed
        calculated_height = base_header_height + (num_data_rows * base_row_height)
        
        # If calculated height exceeds maximum available height, scale down aggressively
        # Use 92% instead of 95% to leave more safety margin
        if calculated_height > max_available_height:
            scale_factor = (max_available_height * 0.92) / calculated_height  # 92% to leave safety margin
            base_row_height *= scale_factor
            base_header_height *= scale_factor
            calculated_height = base_header_height + (num_data_rows * base_row_height)
            print(f"INFO: Scaled down row heights by factor {scale_factor:.2f} to fit {num_data_rows} rows in {max_available_height:.2f} inches")
            
            # If still too tall after scaling, reduce font size further and scale again
            if calculated_height > max_available_height:
                additional_scale = (max_available_height * 0.92) / calculated_height
                base_row_height *= additional_scale
                base_header_height *= additional_scale
                data_font_size = max(6, data_font_size - 1)
                header_font_size = max(7, header_font_size - 1)
                calculated_height = base_header_height + (num_data_rows * base_row_height)
                print(f"INFO: Additional scaling applied: {additional_scale:.2f}, reduced font sizes")
                
                # If STILL too tall, apply even more aggressive scaling
                if calculated_height > max_available_height:
                    emergency_scale = (max_available_height * 0.90) / calculated_height
                    base_row_height *= emergency_scale
                    base_header_height *= emergency_scale
                    calculated_height = base_header_height + (num_data_rows * base_row_height)
                    print(f"WARNING: Emergency scaling applied: {emergency_scale:.2f}")
        
        # Ensure minimum row heights for readability
        base_row_height = max(0.18, base_row_height)
        base_header_height = max(0.22, base_header_height)
        
        # ADAPTIVE FONT SIZE based on available space and row count
        if num_data_rows <= 10:
            data_font_size = 9
            header_font_size = 10
        elif num_data_rows <= 15:
            data_font_size = 8
            header_font_size = 9
        elif num_data_rows <= 20:
            data_font_size = 7
            header_font_size = 8
        elif num_data_rows <= 25:
            data_font_size = 6
            header_font_size = 7
        else:
            data_font_size = 6
            header_font_size = 7
        
        # Reduce font if height is tight (this is checked again after scaling above)
        if calculated_height > max_available_height * 0.85:
            data_font_size = max(6, data_font_size - 1)
            header_font_size = max(7, header_font_size - 1)
        
        # ADAPTIVE COLUMN WIDTH - Ensure all columns fit
        if cols > 0:
            col_width = available_width / cols
            
            # Reduce column width for larger tables
            if num_data_rows > 15:
                col_width *= 0.85
            elif num_data_rows > 10:
                col_width *= 0.90
            
            # Ensure minimum column width
            min_col_width = 0.30
            if col_width < min_col_width:
                col_width = max(min_col_width, available_width / cols * 0.75)
                print(f"WARNING: Column width reduced to {col_width:.2f} for {cols} columns")
            
            # Recalculate to ensure fit
            table_width = min(col_width * cols, available_width)
        else:
            table_width = available_width
        
        # Final table dimensions - ensure table height equals sum of row heights
        # This prevents PowerPoint from stretching the table beyond our calculated dimensions
        # Recalculate exact height from final row heights
        final_calculated_height = base_header_height + (num_data_rows * base_row_height)
        
        # Ensure this doesn't exceed maximum available height (use 92% for safety margin)
        if final_calculated_height > max_available_height:
            # Scale down one more time to ensure fit
            final_scale = (max_available_height * 0.92) / final_calculated_height
            base_row_height *= final_scale
            base_header_height *= final_scale
            final_calculated_height = base_header_height + (num_data_rows * base_row_height)
            print(f"WARNING: Final scaling applied: {final_scale:.2f} to ensure table fits within {max_available_height:.2f} inches")
        
        # Table height must exactly match sum of row heights
        table_height = final_calculated_height
        
        # Final check - ensure table won't overflow slide (leave 0.4" margin for safety)
        table_bottom = top + table_height
        max_allowed_bottom = SLIDE_HEIGHT - 0.4  # Leave 0.4" margin from bottom
        
        if table_bottom > max_allowed_bottom:
            # Reduce even more aggressively to ensure fit
            available_from_top = max_allowed_bottom - top
            if available_from_top > 0.5:  # Only if we have reasonable space
                emergency_scale = (available_from_top * 0.92) / table_height
                base_row_height *= emergency_scale
                base_header_height *= emergency_scale
                table_height = base_header_height + (num_data_rows * base_row_height)
                table_bottom = top + table_height
                print(f"WARNING: Emergency scaling for slide fit: {emergency_scale:.2f}, new height: {table_height:.2f}")
                
                # Verify it now fits - if still too tall, scale again
                if table_bottom > max_allowed_bottom:
                    emergency_scale2 = (available_from_top * 0.90) / table_height
                    base_row_height *= emergency_scale2
                    base_header_height *= emergency_scale2
                    table_height = base_header_height + (num_data_rows * base_row_height)
                    print(f"WARNING: Second emergency scaling: {emergency_scale2:.2f}")
            else:
                print(f"ERROR: Not enough space - top: {top:.2f}, required: {table_height:.2f}, available: {available_from_top:.2f}")
        
        # Create table with calculated dimensions
        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(table_width)
        height_inches = Inches(table_height)
        
        table_shape = slide.shapes.add_table(rows, cols, left_inches, top_inches, width_inches, height_inches)
        table = table_shape.table
        
        # Add table borders for clarity
        # Note: Border setting may not work for all cell types, so we wrap in try-except
        try:
            # Set borders for all cells - light gray borders
            for row in table.rows:
                for cell in row.cells:
                    try:
                        # Try to set borders - not all cells support this
                        if hasattr(cell, 'border_top'):
                            cell.border_top.color.rgb = RGBColor(200, 200, 200)
                            cell.border_top.width = Pt(0.5)
                        if hasattr(cell, 'border_bottom'):
                            cell.border_bottom.color.rgb = RGBColor(200, 200, 200)
                            cell.border_bottom.width = Pt(0.5)
                        if hasattr(cell, 'border_left'):
                            cell.border_left.color.rgb = RGBColor(200, 200, 200)
                            cell.border_left.width = Pt(0.5)
                        if hasattr(cell, 'border_right'):
                            cell.border_right.color.rgb = RGBColor(200, 200, 200)
                            cell.border_right.width = Pt(0.5)
                    except (AttributeError, TypeError) as e:
                        # Some cells don't support border attributes - skip them
                        pass
        except Exception as e:
            # Border setting failed - not critical, continue without borders
            print(f"WARNING: Could not set table borders: {e}")
            pass
        
        # Set column widths
        if cols > 0:
            final_col_width = table_width / cols
            for col_idx in range(cols):
                table.columns[col_idx].width = Inches(final_col_width)
        
        # Add table borders for clarity
        try:
            for row in table.rows:
                for cell in row.cells:
                    try:
                        if hasattr(cell, 'border_top'):
                            cell.border_top.color.rgb = RGBColor(200, 200, 200)
                            cell.border_top.width = Pt(0.5)
                        if hasattr(cell, 'border_bottom'):
                            cell.border_bottom.color.rgb = RGBColor(200, 200, 200)
                            cell.border_bottom.width = Pt(0.5)
                        if hasattr(cell, 'border_left'):
                            cell.border_left.color.rgb = RGBColor(200, 200, 200)
                            cell.border_left.width = Pt(0.5)
                        if hasattr(cell, 'border_right'):
                            cell.border_right.color.rgb = RGBColor(200, 200, 200)
                            cell.border_right.width = Pt(0.5)
                    except (AttributeError, TypeError):
                        pass
        except Exception as e:
            print(f"WARNING: Could not set table borders: {e}")
        
        # Merge formatting with defaults
        if formatting is None:
            formatting = {}
        
        # Default header formatting - use title color #004E6F
        default_header_formatting = {
            "font_size": header_font_size,
            "bold": True,
            "fill_color": "#004E6F",  # Title color shade
            "font_color": "#FFFFFF"   # White text
        }
        
        # Default data formatting - white background, black text
        default_data_formatting = {
            "font_size": data_font_size,
            "bold": False,
            "font_color": "#000000",  # Black text for readability
            "fill_color": "#FFFFFF",  # White background
            "alignment": "left"  # Left align data cells
        }
        
        header_formatting = formatting.get("header_formatting", default_header_formatting)
        data_formatting = formatting.get("data_formatting", default_data_formatting)
        
        # Ensure header formatting has white text
        if "fill_color" in header_formatting:
            header_formatting["font_color"] = "#FFFFFF"
        
        # Populate header row
        for col_idx, col_name in enumerate(data.columns):
            cell = table.cell(0, col_idx)
            # Clean column name for display - handle NaN/None
            if pd.isna(col_name) or col_name is None:
                cell.text = ""
            else:
                cell.text = str(col_name).strip()
            # Enable text wrapping for headers
            cell.text_frame.word_wrap = True
            cell.text_frame.auto_size = None  # Disable auto-size for better control
            self.formatter.format_table_cell(cell, header_formatting)
            # Left align header text (changed from center)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Populate data rows with row type-based formatting
        number_formatting = formatting.get("number_formatting", {})
        
        for row_idx, (_, row_data) in enumerate(data.iterrows(), start=1):
            row_type = row_types[row_idx - 1]  # Get detected row type
            
            # Apply formatting based on row type
            if row_type == 'total':
                # Grand total row: Dark blue background, white text, bold
                row_formatting = {
                    "font_size": data_font_size,
                    "bold": True,
                    "fill_color": "#004E6F",  # Dark blue
                    "font_color": "#FFFFFF",  # White text
                    "alignment": "left"
                }
            elif row_type == 'subtotal':
                # Subtotal row: Light blue background, black text, bold
                row_formatting = {
                    "font_size": data_font_size,
                    "bold": True,
                    "fill_color": "#E6F2F8",  # Light blue
                    "font_color": "#000000",  # Black text
                    "alignment": "left"
                }
            else:
                # Regular row: White background, black text
                row_formatting = data_formatting.copy()
                row_formatting["fill_color"] = "#FFFFFF"
            
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                
                # Clean NaN/None values first
                if pd.isna(value) or value is None or str(value).lower() in ['nan', 'none', 'nat']:
                    cell.text = ""
                else:
                    # Apply number formatting if specified
                    col_name = data.columns[col_idx]
                    if col_name in number_formatting:
                        format_type = number_formatting[col_name]
                        try:
                            value_float = float(value)
                            if format_type == "percentage":
                                # Check if value is already a percentage (>= 1) or decimal (< 1)
                                # Also check if it's a very large number (likely already multiplied)
                                if value_float < 1:
                                    # It's a decimal (0-1), convert to percentage
                                    cell.text = f"{value_float * 100:.1f}%"
                                elif value_float > 100:
                                    # It's been multiplied already, divide by 100
                                    cell.text = f"{value_float / 100:.1f}%"
                                else:
                                    # It's already a percentage (1-100), just add % sign
                                    cell.text = f"{value_float:.1f}%"
                            elif format_type == "percentage_decimal":
                                # Value is already decimal (0-1), convert to percentage
                                cell.text = f"{value_float * 100:.1f}%"
                            elif format_type == "number":
                                # Remove decimal if it's .0
                                if value_float == int(value_float):
                                    cell.text = f"{int(value_float):,}"
                                else:
                                    cell.text = f"{value_float:,.2f}"
                            elif format_type == "integer":
                                cell.text = f"{int(value_float):,}"
                            elif format_type == "currency":
                                cell.text = f"${value_float:,.2f}"
                            else:
                                cell.text = str(value).strip()
                        except (ValueError, TypeError):
                            cell.text = str(value).strip() if value else ""
                    else:
                        # Default formatting - smart percentage detection
                        try:
                            value_float = float(value)
                            if value_float == int(value_float):
                                cell.text = str(int(value_float))
                            else:
                                # Check if this looks like a percentage (between 0-100 and in a percentage column)
                                # Percentage columns are typically after numeric columns
                                if 0 < value_float < 100 and col_idx >= 2:
                                    # Format as percentage if it looks like one
                                    cell.text = f"{value_float:.1f}%"
                                else:
                                    cell.text = str(value).strip()
                        except (ValueError, TypeError):
                            cell.text = str(value).strip() if value else ""
                
                # Apply formatting
                cell_formatting = row_formatting.copy()
                
                # Apply conditional formatting if specified (overrides row color)
                # Use title color shade for conditional formatting
                if "conditional_colors" in formatting:
                    for cond in formatting["conditional_colors"]:
                        if cond.get("column") == col_name:
                            try:
                                value_float = float(value)
                                condition = cond.get("condition", "<")
                                threshold = cond.get("threshold", 0)
                                if (condition == "<" and value_float < threshold) or \
                                   (condition == ">" and value_float > threshold) or \
                                   (condition == "==" and value_float == threshold):
                                    # Use title color shade instead of red
                                    cell_formatting["font_color"] = cond.get("color", "#004E6F")
                            except (ValueError, TypeError):
                                pass
                
                # Enable text wrapping for data cells
                cell.text_frame.word_wrap = True
                cell.text_frame.auto_size = None  # Disable auto-size for better control
                # Set vertical alignment to middle for better appearance
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                self.formatter.format_table_cell(cell, cell_formatting)
                # Left align data
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Apply additional formatting
        if formatting:
            self.formatter.format_table(table, formatting)
        
        # Set row heights with calculated values - MUST match table height exactly
        # Calculate total height from row heights to verify
        total_row_height = base_header_height + (num_data_rows * base_row_height)
        
        # Verify row heights sum matches table height (allow small tolerance)
        if abs(total_row_height - table_height) > 0.01:
            # Adjust row heights proportionally to match table height exactly
            height_ratio = table_height / total_row_height
            base_row_height *= height_ratio
            base_header_height *= height_ratio
            total_row_height = base_header_height + (num_data_rows * base_row_height)
            print(f"INFO: Adjusted row heights by ratio {height_ratio:.3f} to match table height exactly")
        
        # Set row heights
        for row_idx, row in enumerate(table.rows):
            if row_idx == 0:
                row.height = Inches(base_header_height)
            else:
                row.height = Inches(base_row_height)
        
        # Verify final dimensions
        actual_table_height = base_header_height + (num_data_rows * base_row_height)
        table_bottom_pos = top + actual_table_height
        
        print(f"DEBUG: Table created - Rows: {num_data_rows} (Regular: {num_regular_rows}, Subtotal: {num_subtotal_rows}, Total: {num_total_rows})")
        print(f"DEBUG: Font sizes - Header: {header_font_size}, Data: {data_font_size}")
        print(f"DEBUG: Row heights - Header: {base_header_height:.2f}, Data: {base_row_height:.2f}")
        print(f"DEBUG: Table dimensions - Width: {table_width:.2f}, Height: {actual_table_height:.2f} (created with {table_height:.2f})")
        print(f"DEBUG: Table position - Top: {top:.2f}, Bottom: {table_bottom_pos:.2f}, Slide height: {SLIDE_HEIGHT:.2f}")
        print(f"DEBUG: Max available height was: {max_available_height:.2f}")
        
        # Verify table fits
        if table_bottom_pos > SLIDE_HEIGHT:
            print(f"ERROR: Table extends beyond slide! Top: {top:.2f}, Height: {actual_table_height:.2f}, Bottom: {table_bottom_pos:.2f}, Slide: {SLIDE_HEIGHT:.2f}")
        else:
            margin = SLIDE_HEIGHT - table_bottom_pos
            print(f"SUCCESS: Table fits within slide with {margin:.2f} inch margin at bottom")
        
        return table_shape
    
    def add_bullet_list(self, slide, items: List[str], left: float, top: float,
                       width: float, height: float, formatting: Optional[Dict] = None) -> Any:
        """
        Add a bullet list to a slide.
        
        Args:
            slide: Slide object
            items: List of bullet items
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            formatting: Optional formatting dictionary
        
        Returns:
            Shape object
        """
        text_box = self.add_text_box(slide, "", left, top, width, height, formatting)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Add items as paragraphs with bullets
        for i, item in enumerate(items):
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            
            paragraph.text = item
            paragraph.level = 0
            paragraph.space_after = Pt(6)
            
            if formatting:
                self.formatter.format_paragraph(paragraph, formatting)
        
        return text_box
    
    def update_text_in_shape(self, slide, shape_index: int, text: str,
                           formatting: Optional[Dict] = None):
        """
        Update text in an existing shape.
        
        Args:
            slide: Slide object
            shape_index: Index of the shape to update
            text: New text content
            formatting: Optional formatting dictionary
        """
        if shape_index < len(slide.shapes):
            shape = slide.shapes[shape_index]
            if shape.has_text_frame:
                shape.text_frame.text = text
                if formatting:
                    self.formatter.format_text_box(shape.text_frame, formatting)
    
    def update_table_data(self, slide, shape_index: int, data: pd.DataFrame,
                         formatting: Optional[Dict] = None):
        """
        Update data in an existing table.
        
        Args:
            slide: Slide object
            shape_index: Index of the shape to update
            data: New DataFrame data
            formatting: Optional formatting dictionary
        """
        if shape_index < len(slide.shapes):
            shape = slide.shapes[shape_index]
            if shape.has_table:
                table = shape.table
                
                # Clear existing data (except header)
                for row_idx in range(len(table.rows) - 1, 0, -1):
                    table._tbl.remove(table.rows[row_idx]._tr)
                
                # Add new rows
                for _, row_data in data.iterrows():
                    new_row = table.rows.add()
                    for col_idx, value in enumerate(row_data):
                        if col_idx < len(new_row.cells):
                            new_row.cells[col_idx].text = str(value)
                            if formatting and "data_formatting" in formatting:
                                self.formatter.format_table_cell(new_row.cells[col_idx], formatting["data_formatting"])
                
                if formatting:
                    self.formatter.format_table(table, formatting)
    
    def find_shape_by_name(self, slide, name: str) -> Optional[Any]:
        """
        Find a shape by name.
        
        Args:
            slide: Slide object
            name: Name of the shape
        
        Returns:
            Shape object or None
        """
        for shape in slide.shapes:
            if shape.name == name:
                return shape
        return None
    
    def populate_slide_from_mapping(self, slide, data: Dict[str, Any],
                                   mapping: Dict[str, Any]):
        """
        Populate a slide based on a mapping configuration.
        
        Args:
            slide: Slide object
            data: Data dictionary
            mapping: Mapping configuration dictionary
        """
        for shape_mapping in mapping.get("shape_mappings", []):
            shape_index = shape_mapping.get("shape_index")
            mapping_type = shape_mapping.get("mapping_type")
            data_source = shape_mapping.get("data_source")
            
            if mapping_type == "text":
                text_value = self._get_text_value(data, shape_mapping)
                if text_value is not None:
                    self.update_text_in_shape(slide, shape_index, str(text_value))
            
            elif mapping_type == "table":
                table_data = self._get_table_data(data, shape_mapping, return_column_mapping=False)
                # Handle tuple return (DataFrame, mapping) if column_mapping was requested
                if isinstance(table_data, tuple):
                    table_data, _ = table_data
                if table_data is not None:
                    self.update_table_data(slide, shape_index, table_data)
    
    def _get_text_value(self, data: Dict[str, Any], mapping: Dict[str, Any]) -> Optional[str]:
        """Extract text value from data based on mapping."""
        data_source = mapping.get("data_source")
        column = mapping.get("column")
        
        if data_source in data:
            df = data[data_source]
            if isinstance(df, pd.DataFrame) and column in df.columns:
                # Return first value or aggregated value
                aggregate = mapping.get("aggregate", "first")
                if aggregate == "sum":
                    return str(df[column].sum())
                elif aggregate == "mean":
                    return str(df[column].mean())
                else:
                    return str(df[column].iloc[0])
        
        return mapping.get("default_value")
    
    def _get_table_data(self, data: Dict[str, Any], mapping: Dict[str, Any], return_column_mapping: bool = False) -> Optional[pd.DataFrame]:
        """Extract table data from data based on mapping."""
        data_source = mapping.get("data_source")
        sheet_name = mapping.get("sheet")
        
        if not data_source:
            print(f"DEBUG: No data_source specified in mapping: {mapping}")
            return None
        
        print(f"DEBUG: Looking for data_source: '{data_source}' (type: {type(data_source)}), available keys: {list(data.keys())}")
        
        # Normalize data_source for matching (strip whitespace)
        data_source_normalized = str(data_source).strip() if data_source else ""
        
        # Try exact match first
        df_source = None
        if data_source_normalized in data:
            df_source = data[data_source_normalized]
            print(f"DEBUG: Found exact match for data_source: '{data_source_normalized}'")
        else:
            # Try case-insensitive match
            data_source_lower = data_source_normalized.lower()
            for key in data.keys():
                key_normalized = str(key).strip()
                if key_normalized.lower() == data_source_lower:
                    df_source = data[key]
                    print(f"DEBUG: Found case-insensitive match: '{key}' for '{data_source_normalized}'")
                    break
            
            # Try partial match (contains)
            if df_source is None:
                for key in data.keys():
                    key_str = str(key).strip().lower()
                    if data_source_lower in key_str or key_str in data_source_lower:
                        df_source = data[key]
                        print(f"DEBUG: Found partial match: '{key}' for '{data_source_normalized}'")
                        break
        
        if df_source is None:
            print(f"WARNING: Could not find data_source '{data_source}' in available data. Available keys: {list(data.keys())[:10]}")
            # Try to use first available data source as fallback
            if data and len(data) > 0:
                first_key = list(data.keys())[0]
                print(f"WARNING: Using first available data source '{first_key}' as fallback")
                df_source = data[first_key]
            else:
                # Return empty DataFrame with structure instead of None
                print(f"WARNING: No data available. Returning empty DataFrame.")
                return pd.DataFrame({"Message": ["No data available"]})
        
        # Handle nested structure: data[file_name][sheet_name]
        if df_source is not None:
            
            # If it's a dict (multiple sheets), get the specific sheet
            if isinstance(df_source, dict):
                if sheet_name and sheet_name in df_source:
                    df = df_source[sheet_name]
                    print(f"DEBUG: Found exact match for sheet: '{sheet_name}'")
                elif sheet_name:
                    # Try case-insensitive match
                    sheet_lower = str(sheet_name).lower().strip()
                    matched_sheet = None
                    for key in df_source.keys():
                        key_str = str(key).lower().strip()
                        if key_str == sheet_lower:
                            matched_sheet = key
                            break
                    
                    if matched_sheet:
                        df = df_source[matched_sheet]
                        print(f"DEBUG: Found case-insensitive match for sheet: '{matched_sheet}' for '{sheet_name}'")
                    else:
                        # Try partial match
                        for key in df_source.keys():
                            key_str = str(key).lower().strip()
                            if sheet_lower in key_str or key_str in sheet_lower:
                                matched_sheet = key
                                break
                        
                        if matched_sheet:
                            df = df_source[matched_sheet]
                            print(f"DEBUG: Found partial match for sheet: '{matched_sheet}' for '{sheet_name}'")
                        else:
                            print(f"WARNING: Sheet '{sheet_name}' not found in {data_source}. Available sheets: {list(df_source.keys())[:10]}")
                            # Return empty DataFrame instead of None
                            return pd.DataFrame({"Message": [f"Sheet '{sheet_name}' not found"]})
                else:
                    # If no sheet specified, use first sheet
                    if df_source:
                        first_sheet = list(df_source.keys())[0]
                        df = df_source[first_sheet]
                        print(f"DEBUG: No sheet specified, using first sheet: '{first_sheet}'")
                    else:
                        df = None
            elif isinstance(df_source, pd.DataFrame):
                df = df_source
                print(f"DEBUG: Data source is a DataFrame, using directly")
            else:
                print(f"WARNING: Data source '{data_source}' has unsupported type: {type(df_source)}")
                # Return empty DataFrame instead of None
                return pd.DataFrame({"Message": ["Data source type not supported"]})
            
            if df is None or not isinstance(df, pd.DataFrame):
                # Return empty DataFrame instead of None
                print(f"WARNING: Data source returned None or invalid type. Returning empty DataFrame.")
                return pd.DataFrame({"Message": ["No data available"]})
            
            # Apply header row offset if needed (data should already be loaded with correct header)
            # But we can re-read if header_row is specified and different
            header_row = mapping.get("header_row")
            if header_row is not None and header_row != 0:
                # Note: This assumes data was loaded with header=0
                # For now, we'll work with the data as-is
                # In future, we might need to re-read with correct header
                pass
            
            # Apply filters if specified
            filters = mapping.get("filters", [])
            result_df = df.copy()
            
            for filter_def in filters:
                column = filter_def.get("column")
                operator = filter_def.get("operator", "!=")
                value = filter_def.get("value")
                
                if column in result_df.columns:
                    if operator == "!=":
                        if value is None:
                            result_df = result_df[result_df[column].notna()]
                        else:
                            result_df = result_df[result_df[column] != value]
                    elif operator == ">=":
                        result_df = result_df[result_df[column] >= value]
                    elif operator == "<=":
                        result_df = result_df[result_df[column] <= value]
                    elif operator == "==":
                        result_df = result_df[result_df[column] == value]
                    elif operator == "notna":
                        result_df = result_df[result_df[column].notna()]
            
            # Select columns if specified
            columns = mapping.get("columns")
            print(f"DEBUG: Column selection - columns parameter: {columns} (type: {type(columns)}, length: {len(columns) if columns else 'N/A'})")
            print(f"DEBUG: Available columns in DataFrame: {list(result_df.columns)}")
            
            # If columns is specified and not empty, use them (preserve user selection and order)
            if columns is not None and len(columns) > 0:
                print(f"DEBUG: User selected {len(columns)} columns, attempting to match: {columns}")
                available_columns = list(result_df.columns)
                matched_cols = []  # Will preserve user's order
                column_mapping_dict = {}  # Maps user column name to actual column name
                
                columns_to_find = columns if isinstance(columns, list) else [columns]
                
                # Enhanced column matching with multiple strategies
                for user_col in columns_to_find:
                    user_col_str = str(user_col).strip()
                    matched_col = None
                    match_type = None
                    
                    # Strategy 1: Exact match (case-sensitive)
                    if user_col_str in available_columns:
                        matched_col = user_col_str
                        match_type = "exact"
                    
                    # Strategy 2: Exact match (case-insensitive, whitespace-insensitive)
                    if not matched_col:
                        user_col_normalized = user_col_str.lower().strip()
                        for avail_col in available_columns:
                            avail_col_normalized = str(avail_col).strip().lower()
                            if avail_col_normalized == user_col_normalized:
                                matched_col = avail_col  # Use original case
                                match_type = "case_insensitive"
                                break
                    
                    # Strategy 3: Partial match (contains)
                    if not matched_col:
                        user_col_normalized = user_col_str.lower().strip()
                        for avail_col in available_columns:
                            avail_col_normalized = str(avail_col).strip().lower()
                            # Check if one contains the other (both directions)
                            if user_col_normalized in avail_col_normalized or avail_col_normalized in user_col_normalized:
                                # Prefer shorter matches for better accuracy
                                if not matched_col or len(avail_col) < len(matched_col):
                                    matched_col = avail_col
                                    match_type = "partial"
                    
                    # Strategy 4: Fuzzy match (handle common variations)
                    if not matched_col:
                        user_col_normalized = user_col_str.lower().strip()
                        # Remove common prefixes/suffixes and special chars for matching
                        user_col_clean = user_col_normalized.replace('_', ' ').replace('-', ' ').replace('.', ' ')
                        user_col_clean = ' '.join(user_col_clean.split())  # Normalize whitespace
                        
                        for avail_col in available_columns:
                            avail_col_normalized = str(avail_col).strip().lower()
                            avail_col_clean = avail_col_normalized.replace('_', ' ').replace('-', ' ').replace('.', ' ')
                            avail_col_clean = ' '.join(avail_col_clean.split())  # Normalize whitespace
                            
                            if user_col_clean == avail_col_clean:
                                matched_col = avail_col
                                match_type = "fuzzy"
                                break
                    
                    if matched_col:
                        # Only add if not already added (avoid duplicates)
                        if matched_col not in matched_cols:
                            matched_cols.append(matched_col)
                            column_mapping_dict[user_col] = matched_col
                            print(f"INFO: Matched column '{user_col}' -> '{matched_col}' ({match_type})")
                        else:
                            print(f"WARNING: Column '{user_col}' matched to '{matched_col}' which was already matched")
                    else:
                        print(f"WARNING: Could not match column '{user_col}' to any available column")
                        print(f"WARNING:   Available columns: {available_columns[:10]}")
                
                # Preserve user's column order and select matched columns
                if matched_cols:
                    # Reorder matched_cols to preserve user's order
                    ordered_matched_cols = []
                    for user_col in columns_to_find:
                        if user_col in column_mapping_dict:
                            actual_col = column_mapping_dict[user_col]
                            if actual_col not in ordered_matched_cols:
                                ordered_matched_cols.append(actual_col)
                    
                    # Use the ordered list to select columns
                    result_df = result_df[ordered_matched_cols]
                    print(f"INFO: Successfully matched {len(ordered_matched_cols)}/{len(columns_to_find)} columns in user's order")
                    print(f"INFO: Selected columns (in order): {list(result_df.columns)}")
                    print(f"DEBUG: Result DataFrame shape: {result_df.shape}")
                    
                    # Create column mapping if requested
                    if return_column_mapping:
                        return (result_df, column_mapping_dict)
                    
                    return result_df
                else:
                    # No columns matched - this is an error, but return all columns as fallback
                    print(f"ERROR: None of the {len(columns_to_find)} selected columns could be matched!")
                    print(f"ERROR: Selected: {columns_to_find}")
                    print(f"ERROR: Available: {available_columns[:20]}")
                    print(f"INFO: Falling back to all columns to prevent data loss")
                    # Continue to return all columns below
            elif columns is not None and len(columns) == 0:
                # Empty columns list - return all columns (user wants all columns)
                print(f"INFO: Empty columns list - returning all columns")
                # Continue to return all columns below
            else:
                # columns is None - return all columns
                print(f"INFO: No columns specified - returning all columns")
                # Continue to return all columns below
            
            # Limit rows if specified
            max_rows = mapping.get("max_rows")
            if max_rows:
                result_df = result_df.head(max_rows)
            
            # If no column selection was done (columns not specified, empty list, or matching failed), return all columns
            print(f"INFO: Returning DataFrame with all columns. Shape: {result_df.shape}, Columns: {list(result_df.columns)}")
            
            # Validate data integrity - ensure we have data
            if len(result_df) == 0:
                print(f"WARNING: DataFrame is empty (no rows). Data source: {data_source}, Sheet: {sheet_name}")
            if len(result_df.columns) == 0:
                print(f"WARNING: DataFrame has no columns. Data source: {data_source}, Sheet: {sheet_name}")
            
            # Log final result
            print(f"INFO: Final DataFrame - {len(result_df)} rows, {len(result_df.columns)} columns")
            if return_column_mapping:
                # Create identity mapping for all columns
                column_mapping = {col: col for col in result_df.columns}
                print(f"INFO: Column mapping created: {len(column_mapping)} mappings")
                return (result_df, column_mapping)
            
            return result_df
        
        # If we get here, something went wrong - return empty DataFrame instead of None
        print(f"WARNING: _get_table_data reached end without returning data. Returning empty DataFrame.")
        return pd.DataFrame({"Message": ["No data available"]})
    
    def add_chart(self, slide, data: pd.DataFrame, chart_type: str = "column",
                  left: float = 1, top: float = 2, width: float = 8, height: float = 4,
                  x_column: Optional[str] = None, y_columns: Optional[List[str]] = None,
                  title: str = "", formatting: Optional[Dict] = None) -> Any:
        """
        Add a chart to a slide.
        
        Args:
            slide: Slide object
            data: DataFrame containing chart data
            chart_type: Type of chart ("column", "bar", "line", "pie")
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            x_column: Column name for X-axis categories
            y_columns: List of column names for Y-axis values
            title: Chart title
            formatting: Optional formatting dictionary
        
        Returns:
            Shape object
        """
        if data is None or len(data) == 0:
            raise ValueError("Cannot create chart with empty data")
        
        print(f"DEBUG: Creating chart with data shape: {data.shape}, columns: {list(data.columns)}")
        
        # Default to first column as X, remaining as Y
        if x_column is None:
            x_column = data.columns[0] if len(data.columns) > 0 else None
        
        if y_columns is None or len(y_columns) == 0:
            y_columns = [col for col in data.columns if col != x_column]
            if len(y_columns) == 0:
                y_columns = [data.columns[0]] if len(data.columns) > 0 else []
        
        if not x_column or not y_columns:
            raise ValueError(f"Cannot create chart: x_column={x_column}, y_columns={y_columns}")
        
        # Verify columns exist in data
        if x_column not in data.columns:
            # Try to find matching column
            x_col_lower = str(x_column).lower().strip()
            for col in data.columns:
                if str(col).lower().strip() == x_col_lower:
                    x_column = col
                    break
            if x_column not in data.columns:
                raise ValueError(f"X column '{x_column}' not found in data. Available: {list(data.columns)}")
        
        # Filter y_columns to only those that exist
        valid_y_columns = []
        for y_col in (y_columns if isinstance(y_columns, list) else [y_columns]):
            if y_col in data.columns:
                valid_y_columns.append(y_col)
            else:
                # Try case-insensitive match
                y_col_lower = str(y_col).lower().strip()
                for col in data.columns:
                    if str(col).lower().strip() == y_col_lower:
                        valid_y_columns.append(col)
                        print(f"DEBUG: Matched y_column '{y_col}' to '{col}'")
                        break
        
        if not valid_y_columns:
            raise ValueError(f"No valid Y columns found. Requested: {y_columns}, Available: {list(data.columns)}")
        
        y_columns = valid_y_columns
        print(f"DEBUG: Using x_column='{x_column}', y_columns={y_columns}")
        
        # Prepare chart data
        chart_data = CategoryChartData()
        
        # Get categories from X column
        categories = data[x_column].astype(str).tolist()
        chart_data.categories = categories
        print(f"DEBUG: Categories ({len(categories)}): {categories[:5]}...")
        
        # Add series for each Y column
        for y_col in y_columns:
            if y_col in data.columns and y_col != x_column:
                values = data[y_col].tolist()
                # Convert to numeric values
                numeric_values = []
                for val in values:
                    try:
                        # Handle percentage strings, commas, etc.
                        val_str = str(val).replace('%', '').replace(',', '').strip()
                        numeric_values.append(float(val_str))
                    except (ValueError, TypeError):
                        numeric_values.append(0.0)
                
                print(f"DEBUG: Added series '{y_col}' with {len(numeric_values)} values: {numeric_values[:5]}...")
                chart_data.add_series(y_col, numeric_values)
        
        # Map chart type string to enum
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "area": XL_CHART_TYPE.AREA,
        }
        
        chart_type_enum = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # Add chart to slide
        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(width)
        height_inches = Inches(height)
        
        graphic_frame = slide.shapes.add_chart(
            chart_type_enum, left_inches, top_inches, width_inches, height_inches, chart_data
        )
        chart = graphic_frame.chart
        
        # Set chart title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            # Format chart title
            for paragraph in chart.chart_title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)
                    run.font.bold = True
        
        # Configure legend - enable and position at bottom for better alignment
        chart.has_legend = True
        try:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            # Format legend font
            chart.legend.font.size = Pt(10)
        except (AttributeError, ValueError) as e:
            print(f"WARNING: Could not configure legend: {e}")
        
        # Configure axis titles and formatting
        try:
            # Set X-axis (category axis) title
            chart.category_axis.has_title = True
            x_axis_title = formatting.get("x_axis_title") if formatting else None
            if x_axis_title:
                chart.category_axis.axis_title.text_frame.text = str(x_axis_title)
            else:
                # Use x_column name as X-axis title
                chart.category_axis.axis_title.text_frame.text = str(x_column)
            
            # Format X-axis title with title color
            title_color = "#004E6F"  # Title font color for visibility on white background
            if title_color.startswith("#"):
                title_color = title_color[1:]
            r_title = int(title_color[0:2], 16)
            g_title = int(title_color[2:4], 16)
            b_title = int(title_color[4:6], 16)
            
            for paragraph in chart.category_axis.axis_title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.bold = False
                    run.font.color.rgb = RGBColor(r_title, g_title, b_title)
            
            # Format X-axis labels (category labels) with title color for visibility
            chart.category_axis.tick_labels.font.size = Pt(9)
            chart.category_axis.tick_labels.font.color.rgb = RGBColor(r_title, g_title, b_title)
            
            # Set Y-axis (value axis) title
            chart.value_axis.has_title = True
            y_axis_title = formatting.get("y_axis_title") if formatting else None
            if y_axis_title:
                chart.value_axis.axis_title.text_frame.text = str(y_axis_title)
            else:
                # Use first Y column name or default title
                if len(y_columns) == 1:
                    chart.value_axis.axis_title.text_frame.text = str(y_columns[0])
                else:
                    chart.value_axis.axis_title.text_frame.text = "Value"
            
            # Format Y-axis title with title color
            for paragraph in chart.value_axis.axis_title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.bold = False
                    run.font.color.rgb = RGBColor(r_title, g_title, b_title)
            
            # Format Y-axis labels (value labels) with title color for visibility
            chart.value_axis.tick_labels.font.size = Pt(9)
            chart.value_axis.tick_labels.font.color.rgb = RGBColor(r_title, g_title, b_title)
            
            # Configure gridlines for better readability and insights
            # Enable major gridlines on value axis (Y-axis)
            try:
                chart.value_axis.has_major_gridlines = True
                major_gridlines = chart.value_axis.major_gridlines
                major_gridlines.format.line.color.rgb = RGBColor(180, 180, 180)  # Light gray for visibility
                major_gridlines.format.line.width = Pt(0.75)  # Visible but not too thick
                print(f"DEBUG: Enabled major gridlines on Y-axis")
            except (AttributeError, ValueError) as e:
                print(f"WARNING: Could not configure major gridlines: {e}")
            
            # Enable minor gridlines for more granular reading (optional, lighter)
            try:
                chart.value_axis.has_minor_gridlines = True
                minor_gridlines = chart.value_axis.minor_gridlines
                if minor_gridlines:
                    minor_gridlines.format.line.color.rgb = RGBColor(220, 220, 220)  # Very light gray
                    minor_gridlines.format.line.width = Pt(0.5)  # Thinner than major
                    print(f"DEBUG: Enabled minor gridlines on Y-axis")
            except (AttributeError, ValueError) as e:
                print(f"WARNING: Could not configure minor gridlines: {e}")
            
            # For column/bar charts, also enable vertical gridlines on category axis if helpful
            # (Usually not needed for category axis, but can be enabled if data is dense)
            try:
                chart.category_axis.has_major_gridlines = True
                cat_major_gridlines = chart.category_axis.major_gridlines
                if cat_major_gridlines:
                    cat_major_gridlines.format.line.color.rgb = RGBColor(200, 200, 200)  # Light gray
                    cat_major_gridlines.format.line.width = Pt(0.5)  # Thinner vertical lines
                    print(f"DEBUG: Enabled major gridlines on X-axis")
            except (AttributeError, ValueError) as e:
                # Category axis gridlines may not be available for all chart types
                pass
            
        except (AttributeError, ValueError) as e:
            print(f"WARNING: Could not configure axes: {e}")
        
        # Apply formatting - use title color shades (#004E6F) for charts
        title_color_base = "#004E6F"  # Title color
        # Generate shades of title color for multiple series
        def get_color_shade(base_hex, index, total):
            """Generate a shade of the base color."""
            if base_hex.startswith("#"):
                base_hex = base_hex[1:]
            r_base = int(base_hex[0:2], 16)
            g_base = int(base_hex[2:4], 16)
            b_base = int(base_hex[4:6], 16)
            
            # Create lighter shades by mixing with white
            # First series: darkest (original), subsequent: progressively lighter
            if total == 1:
                return (r_base, g_base, b_base)
            
            # Mix with white: (1 - factor) * base + factor * 255
            factor = index / (total - 1) * 0.4  # Max 40% lighter
            r = int((1 - factor) * r_base + factor * 255)
            g = int((1 - factor) * g_base + factor * 255)
            b = int((1 - factor) * b_base + factor * 255)
            return (r, g, b)
        
        # Apply colors to chart series - use title color shades
        num_series = len(chart.series)
        for i, series in enumerate(chart.series):
            # Use title color shades if colors not specified, or use provided colors
            if formatting and "colors" in formatting and len(formatting.get("colors", [])) > i:
                color_str = formatting["colors"][i]
                if color_str.startswith("#"):
                    color_str = color_str[1:]
                try:
                    r, g, b = int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16)
                except (ValueError, IndexError):
                    r, g, b = get_color_shade(title_color_base, i, num_series)
            else:
                # Default: use title color shades
                r, g, b = get_color_shade(title_color_base, i, num_series)
            
            try:
                # For column/bar charts, use fill color
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(r, g, b)
            except (ValueError, AttributeError):
                pass
            
            # For line charts, enhance line visibility and add markers
            try:
                # Check if this is a line or area chart type
                is_line_chart = chart_type.lower() in ["line", "area"]
                # Also check chart type enum for robustness
                if hasattr(chart, 'chart_type'):
                    chart_type_enum_str = str(chart.chart_type)
                    is_line_chart = is_line_chart or 'LINE' in chart_type_enum_str or 'AREA' in chart_type_enum_str
                
                if is_line_chart:
                    line = series.line
                    line.color.rgb = RGBColor(r, g, b)
                    line.width = Pt(2.5)  # Thicker lines for better visibility
                    
                    # Add markers to line charts for better insights
                    if chart_type.lower() == "line":
                        try:
                            marker = series.marker
                            if marker:
                                marker.style = None  # Use default marker style
                                marker.size = 6  # Marker size
                                marker.format.fill.solid()
                                marker.format.fill.fore_color.rgb = RGBColor(r, g, b)
                                marker.format.line.color.rgb = RGBColor(r, g, b)
                        except (AttributeError, ValueError):
                            # Markers might not be available for all line chart types
                            pass
            except (AttributeError, ValueError) as e:
                # Some chart types don't support line properties
                pass
        
        return graphic_frame

