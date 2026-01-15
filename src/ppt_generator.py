"""
PPT Generator
Main PowerPoint generation orchestrator.
"""

import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, List, Any, Optional
import yaml
import pandas as pd
try:
    from .ppt_builder import PPTBuilder
    from .ppt_formatter import PPTFormatter
except ImportError:
    from ppt_builder import PPTBuilder
    from ppt_formatter import PPTFormatter


class PPTGenerator:
    """Generates PowerPoint decks from processed data."""
    
    @staticmethod
    def _clean_text(text: Any) -> str:
        """Clean text by removing NaN/None values and trimming whitespace."""
        import pandas as pd
        if text is None or pd.isna(text) or str(text).lower() in ['nan', 'none', 'nat']:
            return ""
        return str(text).strip()
    
    def __init__(self, template_path: Optional[str] = None,
                 slides_config: Optional[str] = None,
                 formatting_config: Optional[str] = None,
                 affiliate: Optional[str] = None):
        """
        Initialize the PPT generator.
        
        Args:
            template_path: Path to PowerPoint template file
            slides_config: Path to slides configuration YAML file
            formatting_config: Path to formatting configuration YAML file
            affiliate: Selected affiliate (AIL, APC, ASC) for replacing in title slide
        """
        self.template_path = template_path
        self.slides_config = slides_config
        self.formatting_config = formatting_config
        self.affiliate = affiliate
        
        # Load template
        if template_path and os.path.exists(template_path):
            self.presentation = Presentation(template_path)
        else:
            self.presentation = Presentation()
        
        # Load configurations
        self.slides_mapping = {}
        self.formatting_rules = {}
        
        if slides_config and os.path.exists(slides_config):
            self._load_slides_config()
        
        if formatting_config and os.path.exists(formatting_config):
            self._load_formatting_config()
        
        # Initialize formatter and builder
        self.formatter = PPTFormatter(self.formatting_rules)
        self.builder = PPTBuilder(self.presentation, self.formatter)
    
    def _load_slides_config(self):
        """Load slides configuration from YAML file."""
        with open(self.slides_config, 'r', encoding='utf-8') as f:
            config = yaml.safe_load(f)
            self.slides_mapping = config.get("slides", [])
            print(f"DEBUG: Loaded {len(self.slides_mapping)} slide configurations")
            for idx, slide_config in enumerate(self.slides_mapping, start=1):
                print(f"DEBUG: Slide {idx}: type={slide_config.get('slide_type')}, title='{slide_config.get('title')}', chart_enabled={slide_config.get('chart', {}).get('enabled', False)}")
    
    def _load_formatting_config(self):
        """Load formatting configuration from YAML file."""
        with open(self.formatting_config, 'r', encoding='utf-8') as f:
            self.formatting_rules = yaml.safe_load(f)
    
    def generate(self, data: Dict[str, Any], output_path: str):
        """
        Generate PowerPoint deck from data.
        
        Args:
            data: Dictionary mapping data source names to DataFrames
            output_path: Path to save the generated PowerPoint file
        """
        # Check if template has slides to preserve
        template_slide_count = len(self.presentation.slides)
        preserve_template = template_slide_count >= 2
        
        if preserve_template:
            # Template has title slide (index 0) and end slide (index -1)
            # We need to preserve these and insert generated slides between them
            
            # Replace affiliate in title slide
            if self.affiliate:
                self._replace_affiliate_in_title_slide(self.affiliate)
            
            # Save end slide information before removing
            end_slide = self.presentation.slides[-1]
            end_slide_layout = end_slide.slide_layout
            end_slide_shapes_info = []
            
            # Extract end slide shapes information for later recreation
            for shape in end_slide.shapes:
                shape_info = {
                    'layout': end_slide_layout,
                    'shape_type': shape.shape_type if hasattr(shape, 'shape_type') else None
                }
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    shape_info['text'] = shape.text_frame.text
                end_slide_shapes_info.append(shape_info)
            
            # Remove all slides except the first one (title slide)
            # Remove slides from end to beginning to avoid index shifting
            while len(self.presentation.slides) > 1:
                try:
                    rId = self.presentation.slides._sldIdLst[-1].rId
                    self.presentation.part.drop_rel(rId)
                    del self.presentation.slides._sldIdLst[-1]
                except (IndexError, AttributeError):
                    break
            
            # Generate slides based on configuration
            print(f"DEBUG: Generating {len(self.slides_mapping)} slides from configuration")
            for idx, slide_config in enumerate(self.slides_mapping, start=1):
                try:
                    print(f"DEBUG: Processing slide {idx} of {len(self.slides_mapping)}")
                    self._generate_slide(slide_config, data)
                    print(f"DEBUG: Successfully completed slide {idx}")
                except Exception as e:
                    import traceback
                    error_msg = f"Failed to generate slide {idx}: {str(e)}\n{traceback.format_exc()}"
                    print(f"ERROR: {error_msg}")
                    # Continue with next slide instead of stopping
                    continue
            
            # Re-add end slide at the end by reloading template and copying the last slide
            try:
                if self.template_path and os.path.exists(self.template_path):
                    template_prs = Presentation(self.template_path)
                    if len(template_prs.slides) >= 2:
                        # Get the end slide from template
                        template_end_slide = template_prs.slides[-1]
                        # Add it to our presentation using the same layout
                        end_slide_new = self.presentation.slides.add_slide(template_end_slide.slide_layout)
                        # Copy text from template end slide shapes
                        for i, shape in enumerate(template_end_slide.shapes):
                            if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text:
                                # Try to find corresponding shape in new slide and update text
                                if i < len(end_slide_new.shapes):
                                    new_shape = end_slide_new.shapes[i]
                                    if hasattr(new_shape, 'text_frame') and new_shape.text_frame:
                                        new_shape.text_frame.text = shape.text_frame.text
                        print(f"DEBUG: Re-added end slide from template")
            except Exception as e:
                import traceback
                print(f"WARNING: Could not re-add end slide: {e}\n{traceback.format_exc()}")
                # Fallback: just add slide with layout
                try:
                    self.presentation.slides.add_slide(end_slide_layout)
                    print(f"DEBUG: Added end slide with layout only (content may be missing)")
                except:
                    pass
        else:
            # No template slides to preserve, clear all and generate fresh
            while len(self.presentation.slides) > 0:
                try:
                    rId = self.presentation.slides._sldIdLst[0].rId
                    self.presentation.part.drop_rel(rId)
                    del self.presentation.slides._sldIdLst[0]
                except (IndexError, AttributeError):
                    break
            
            # Generate slides based on configuration
            print(f"DEBUG: Generating {len(self.slides_mapping)} slides from configuration")
            for idx, slide_config in enumerate(self.slides_mapping, start=1):
                try:
                    print(f"DEBUG: Processing slide {idx} of {len(self.slides_mapping)}")
                    self._generate_slide(slide_config, data)
                    print(f"DEBUG: Successfully completed slide {idx}")
                except Exception as e:
                    import traceback
                    error_msg = f"Failed to generate slide {idx}: {str(e)}\n{traceback.format_exc()}"
                    print(f"ERROR: {error_msg}")
                    # Continue with next slide instead of stopping
                    continue
        
        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
        
        # Save presentation
        self.presentation.save(output_path)
        print(f"PowerPoint deck saved to: {output_path}")
    
    def _replace_affiliate_in_title_slide(self, affiliate: str):
        """Replace AIL text with selected affiliate in title slide."""
        if len(self.presentation.slides) == 0:
            return
        
        title_slide = self.presentation.slides[0]
        for shape in title_slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace("AIL", affiliate)
    
    def _set_slide_background(self, slide, color_hex: str = "#FFFFFF"):
        """Set slide background color."""
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            if color_hex.startswith("#"):
                color_hex = color_hex[1:]
            r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
            fill.fore_color.rgb = RGBColor(r, g, b)
        except Exception as e:
            print(f"WARNING: Could not set slide background: {e}")
    
    def _clear_placeholder_text(self, slide):
        """Clear all placeholder text from template shapes (especially yellow eyebrow text)."""
        try:
            for shape in slide.shapes:
                # Skip tables - we'll handle those separately
                if hasattr(shape, 'has_table') and shape.has_table:
                    continue
                
                # Check if shape has text frame
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    try:
                        text_frame = shape.text_frame
                        
                        # Get shape name for debugging
                        shape_name = getattr(shape, 'name', 'Unknown')
                        shape_type = getattr(shape, 'shape_type', 'Unknown')
                        
                        # Get text content
                        text_content = ""
                        try:
                            text_content = text_frame.text.strip() if text_frame.text else ""
                        except:
                            pass
                        
                        # Check if this looks like placeholder/eyebrow text
                        is_placeholder = False
                        
                        # Check text color - if it's yellow or a light yellow color, likely placeholder
                        try:
                            if text_frame.paragraphs:
                                for paragraph in text_frame.paragraphs:
                                    if paragraph.runs:
                                        for run in paragraph.runs:
                                            try:
                                                if hasattr(run, 'font') and hasattr(run.font, 'color'):
                                                    if hasattr(run.font.color, 'rgb'):
                                                        rgb = run.font.color.rgb
                                                        if rgb:
                                                            # Extract RGB values
                                                            r = getattr(rgb, 'r', 0) if hasattr(rgb, 'r') else (rgb[0] if isinstance(rgb, (list, tuple)) else 0)
                                                            g = getattr(rgb, 'g', 0) if hasattr(rgb, 'g') else (rgb[1] if isinstance(rgb, (list, tuple)) else 0)
                                                            b = getattr(rgb, 'b', 0) if hasattr(rgb, 'b') else (rgb[2] if isinstance(rgb, (list, tuple)) else 0)
                                                            
                                                            # Yellow or light yellow detection
                                                            # Yellow is typically RGB(255, 255, 0) or variations
                                                            # Light yellow can be RGB(255, 255, 200) or similar
                                                            if (r > 200 and g > 200 and b < 150):
                                                                is_placeholder = True
                                                                break
                                            except:
                                                pass
                        except:
                            pass
                        
                        # Check for common placeholder text patterns (case-insensitive)
                        if text_content:
                            placeholder_patterns = [
                                'EYEBROW', 'EYEBROW IDENTIFICATION', 'Click to add notes',
                                'Click to edit', 'Placeholder', 'Sample text', 'CALIBRI',
                                '14PT REGULAR', '14PT', 'REGULAR'
                            ]
                            
                            text_upper = text_content.upper()
                            if any(pattern in text_upper for pattern in placeholder_patterns):
                                is_placeholder = True
                            
                            # Also check if text is very short and seems like formatting instructions
                            if len(text_content) < 50 and any(word in text_upper for word in ['PT', 'FONT', 'COLOR', 'REGULAR', 'BOLD']):
                                is_placeholder = True
                        
                        # Clear placeholder text
                        if is_placeholder and text_content:
                            try:
                                # Clear all text from this shape
                                text_frame.clear()
                                # Ensure it's really empty
                                for paragraph in text_frame.paragraphs:
                                    paragraph.clear()
                                    for run in paragraph.runs[:]:
                                        run.text = ""
                                
                                print(f"DEBUG: Cleared placeholder text from shape '{shape_name}' (type: {shape_type}): '{text_content[:50]}'")
                            except Exception as e:
                                print(f"WARNING: Could not clear text from shape '{shape_name}': {e}")
                    except Exception as e:
                        # Skip shapes that cause errors
                        pass
                            
        except Exception as e:
            print(f"WARNING: Error clearing placeholder text: {e}")
    
    def _generate_slide(self, slide_config: Dict, data: Dict[str, Any]):
        """Generate a single slide based on configuration."""
        slide_number = slide_config.get("slide_number", 1)
        slide_type = slide_config.get("slide_type", "content")
        
        print(f"DEBUG: Generating slide {slide_number} of type '{slide_type}'")
        
        try:
            # Get slide layout
            layout = None
            if self.template_path:
                layout_name = slide_config.get("layout_name")
                if layout_name:
                    for layout_option in self.presentation.slide_layouts:
                        if layout_option.name == layout_name:
                            layout = layout_option
                            break
            
            # Add slide
            slide = self.builder.add_slide(layout)
            print(f"DEBUG: Slide {slide_number} created successfully")
            
            # Clear all placeholder text (especially yellow eyebrow text) from template shapes
            self._clear_placeholder_text(slide)
            
            # Set white background for generated slides (not title/end slides from template)
            # Check if this is a generated slide (not template slide)
            current_slide_index = len(self.presentation.slides) - 1
            # Title slide is at index 0, end slide will be at the end
            # All slides in between are generated slides
            if current_slide_index > 0:  # Not the first slide (title slide)
                self._set_slide_background(slide, "#FFFFFF")
            
            # Populate slide based on type
            if slide_type == "title":
                self._generate_title_slide(slide, slide_config, data)
            elif slide_type == "content":
                self._generate_content_slide(slide, slide_config, data, slide_number)
            elif slide_type == "table":
                self._generate_table_slide(slide, slide_config, data, slide_number)
            elif slide_type == "bullet_list":
                self._generate_bullet_list_slide(slide, slide_config, data)
            else:
                # Generic slide generation
                self._generate_generic_slide(slide, slide_config, data)
            
            print(f"DEBUG: Slide {slide_number} populated successfully")
            
        except Exception as e:
            import traceback
            error_msg = f"Error generating slide {slide_number}: {str(e)}\n{traceback.format_exc()}"
            print(f"ERROR: {error_msg}")
            # Don't add error messages to slides - create professional empty state instead
            try:
                if 'slide' not in locals():
                    slide = self.builder.add_slide(None)
                
                # Add title if we have it
                title = slide_config.get("title", f"Slide {slide_number}")
                if title:
                    try:
                        self.builder.add_text_box(
                            slide, str(title),
                            left=0.5, top=0.5, width=9, height=0.8,
                            formatting={"font_size": 14, "font_color": "#004E6F", "font_name": "Calibri", "bold": True, "alignment": "left"}
                        )
                    except:
                        pass
                
                # Add professional "Content unavailable" message instead of error
                try:
                    self.builder.add_text_box(
                        slide, "Content unavailable - please check data configuration",
                        left=2, top=3.5, width=6, height=1.5,
                        formatting={"font_size": 12, "font_color": "#666666", "alignment": "center"}
                    )
                except:
                    pass
            except Exception as inner_e:
                print(f"ERROR: Could not create fallback slide: {inner_e}")
                pass  # If even fallback fails, continue
    
    def _generate_title_slide(self, slide, slide_config: Dict, data: Dict[str, Any]):
        """Generate a title slide."""
        title = slide_config.get("title", "Title")
        subtitle = slide_config.get("subtitle", "")
        
        # Get formatting from config
        title_formatting = slide_config.get("title_formatting", {})
        subtitle_formatting = slide_config.get("subtitle_formatting", {})
        
        # Default title formatting with colors and sizes
        default_title_formatting = {
            "font_size": title_formatting.get("font_size", 48),
            "bold": title_formatting.get("bold", True),
            "font_color": title_formatting.get("font_color", "#003B55"),  # Dark blue
            "alignment": "center",
            "font_name": title_formatting.get("font_name", "Calibri")
        }
        
        # Default subtitle formatting
        default_subtitle_formatting = {
            "font_size": subtitle_formatting.get("font_size", 28),
            "bold": subtitle_formatting.get("bold", False),
            "font_color": subtitle_formatting.get("font_color", "#666666"),  # Gray
            "alignment": "center",
            "font_name": subtitle_formatting.get("font_name", "Calibri")
        }
        
        # Merge with provided formatting
        title_formatting = {**default_title_formatting, **title_formatting}
        subtitle_formatting = {**default_subtitle_formatting, **subtitle_formatting}
        
        # Get title from data if specified
        if "title_data_source" in slide_config:
            title = self._get_text_from_data(data, slide_config["title_data_source"])
        
        # Get subtitle from data if specified
        if "subtitle_data_source" in slide_config:
            subtitle = self._get_text_from_data(data, slide_config["subtitle_data_source"])
        
        # Try to update existing placeholders first
        title_found = False
        subtitle_found = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_name = shape.name.lower()
                if "title" in shape_name or (shape_name == "" and not title_found):
                    shape.text_frame.text = title
                    # Format title with colors
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(title_formatting["font_size"])
                            run.font.bold = title_formatting["bold"]
                            # Apply color
                            if "font_color" in title_formatting:
                                color_str = title_formatting["font_color"]
                                if color_str.startswith("#"):
                                    color_str = color_str[1:]
                                r, g, b = int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16)
                                from pptx.dml.color import RGBColor
                                run.font.color.rgb = RGBColor(r, g, b)
                    title_found = True
                elif "subtitle" in shape_name:
                    shape.text_frame.text = subtitle
                    # Format subtitle with colors
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(subtitle_formatting["font_size"])
                            run.font.bold = subtitle_formatting["bold"]
                            # Apply color
                            if "font_color" in subtitle_formatting:
                                color_str = subtitle_formatting["font_color"]
                                if color_str.startswith("#"):
                                    color_str = color_str[1:]
                                r, g, b = int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16)
                                from pptx.dml.color import RGBColor
                                run.font.color.rgb = RGBColor(r, g, b)
                    subtitle_found = True
        
        # If no placeholders found, add text boxes manually
        if not title_found and title:
            self.builder.add_text_box(
                slide, title, 
                left=1, top=3, width=8, height=1,
                formatting=title_formatting
            )
        if not subtitle_found and subtitle:
            self.builder.add_text_box(
                slide, subtitle,
                left=1, top=4.5, width=8, height=0.8,
                formatting=subtitle_formatting
            )
    
    def _generate_content_slide(self, slide, slide_config: Dict, data: Dict[str, Any], slide_number: int = 1):
        """Generate a content slide."""
        title = slide_config.get("title", "")
        subtitle = slide_config.get("subtitle", "")
        content_mappings = slide_config.get("content_mappings", [])
        chart_config = slide_config.get("chart", None)
        
        # Get formatting from config
        title_formatting = slide_config.get("title_formatting", {})
        subtitle_formatting = slide_config.get("subtitle_formatting", {})
        
        # Default title formatting - use defaults only if user hasn't specified
        default_title_formatting = {
            "font_size": 14,  # Default: 14
            "bold": True,
            "font_color": "#004E6F",  # Default: #004E6F
            "font_name": "Calibri",  # Default: Calibri
            "alignment": "left"  # Left align
        }
        
        # Default subtitle formatting - use defaults only if user hasn't specified
        default_subtitle_formatting = {
            "font_size": 18,  # Default: 18
            "bold": False,
            "font_color": "#009CDE",  # Default: #009CDE
            "font_name": "Georgia",  # Default: Georgia
            "alignment": "left"  # Left align
        }
        
        # Merge: user's formatting overrides defaults
        title_formatting = {**default_title_formatting, **title_formatting}
        subtitle_formatting = {**default_subtitle_formatting, **subtitle_formatting}
        
        # Clean title and subtitle text
        title = self._clean_text(title)
        subtitle = self._clean_text(subtitle)
        
        # Try to set title in existing placeholder
        title_found = False
        if title:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_name = shape.name.lower()
                    if "title" in shape_name or (shape_name == "" and not title_found):
                        shape.text_frame.text = title
                        # Enable word wrap for title
                        shape.text_frame.word_wrap = True
                        # Format title with colors, font, and alignment
                        for paragraph in shape.text_frame.paragraphs:
                            # Set alignment
                            if title_formatting.get("alignment") == "left":
                                paragraph.alignment = PP_ALIGN.LEFT
                            for run in paragraph.runs:
                                # Set font size
                                run.font.size = Pt(title_formatting.get("font_size", 14))
                                # Set font name
                                if "font_name" in title_formatting:
                                    run.font.name = title_formatting["font_name"]
                                # Set bold
                                run.font.bold = title_formatting.get("bold", True)
                                # Apply color
                                if "font_color" in title_formatting:
                                    color_str = title_formatting["font_color"]
                                    if color_str.startswith("#"):
                                        color_str = color_str[1:]
                                    r, g, b = int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16)
                                    from pptx.dml.color import RGBColor
                                    run.font.color.rgb = RGBColor(r, g, b)
                        title_found = True
                        break
        
        # Try to set subtitle in existing placeholder
        subtitle_found = False
        if subtitle:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_name = shape.name.lower()
                    if "subtitle" in shape_name and not subtitle_found:
                        shape.text_frame.text = subtitle
                        # Enable word wrap for subtitle
                        shape.text_frame.word_wrap = True
                        # Format subtitle with colors, font, and alignment
                        for paragraph in shape.text_frame.paragraphs:
                            # Set alignment
                            if subtitle_formatting.get("alignment") == "left":
                                paragraph.alignment = PP_ALIGN.LEFT
                            for run in paragraph.runs:
                                # Set font size
                                run.font.size = Pt(subtitle_formatting.get("font_size", 18))
                                # Set font name
                                if "font_name" in subtitle_formatting:
                                    run.font.name = subtitle_formatting["font_name"]
                                # Set bold
                                run.font.bold = subtitle_formatting.get("bold", False)
                                # Apply color
                                if "font_color" in subtitle_formatting:
                                    color_str = subtitle_formatting["font_color"]
                                    if color_str.startswith("#"):
                                        color_str = color_str[1:]
                                    r, g, b = int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16)
                                    from pptx.dml.color import RGBColor
                                    run.font.color.rgb = RGBColor(r, g, b)
                        subtitle_found = True
                        break
        
        # If no placeholder found, add title as text box
        if not title_found and title:
            # Ensure alignment is set
            title_formatting["alignment"] = "left"
            self.builder.add_text_box(
                slide, title,
                left=0.5, top=0.5, width=9, height=1,
                formatting=title_formatting
            )
        
        # Add subtitle if present and not found in placeholder
        if subtitle and not subtitle_found:
            # Ensure alignment is set
            subtitle_formatting["alignment"] = "left"
            self.builder.add_text_box(
                slide, subtitle,
                left=0.5, top=1.5, width=9, height=0.6,
                formatting=subtitle_formatting
            )
        
        # Track if chart was successfully added
        chart_success = False
        
        # Add chart if configured
        top_offset = 2.5 if subtitle else 1.8
        if chart_config and chart_config.get("enabled", False):
            print(f"DEBUG: Chart is enabled for slide {slide_number}")
            # Build proper chart mapping with all necessary columns (x + y)
            x_column = chart_config.get("x_column")
            y_columns = chart_config.get("y_columns", [])
            
            if x_column and y_columns:
                # Use same data source as content mappings if not specified in chart config
                chart_data_source = chart_config.get("data_source")
                chart_sheet = chart_config.get("sheet")
                chart_header_row = chart_config.get("header_row", 0)
                
                # Create a mapping for chart data that includes all needed columns
                chart_mapping = {
                    "data_source": chart_data_source,
                    "sheet": chart_sheet,
                    "header_row": chart_header_row,
                    "columns": [x_column] + (y_columns if isinstance(y_columns, list) else [y_columns])
                }
                
                print(f"DEBUG: Chart mapping - data_source: {chart_mapping['data_source']}, sheet: {chart_mapping['sheet']}")
                print(f"DEBUG: Chart mapping - columns: {chart_mapping['columns']}")
                
                try:
                    chart_data = self.builder._get_table_data(data, chart_mapping, return_column_mapping=True)
                    
                    if chart_data is not None:
                        if isinstance(chart_data, tuple):
                            # Returned (DataFrame, column_mapping)
                            chart_data_df, column_mapping = chart_data
                            chart_data = chart_data_df
                        else:
                            # Just DataFrame, create mapping from column order
                            column_mapping = {}
                            actual_columns = list(chart_data.columns)
                            requested_columns = [x_column] + (y_columns if isinstance(y_columns, list) else [y_columns])
                            for i, req_col in enumerate(requested_columns):
                                if i < len(actual_columns):
                                    column_mapping[req_col] = actual_columns[i]
                        
                        print(f"DEBUG: Chart data retrieved - shape: {chart_data.shape}, columns: {list(chart_data.columns)}")
                        print(f"DEBUG: Column mapping: {column_mapping}")
                        
                        if len(chart_data) > 0 and len(chart_data.columns) > 0:
                            # Use column mapping to find actual X and Y column names
                            actual_x_column = column_mapping.get(x_column)
                            if not actual_x_column:
                                # Try to find it in the DataFrame directly
                                x_col_normalized = str(x_column).strip().lower()
                                for col in chart_data.columns:
                                    if str(col).strip().lower() == x_col_normalized:
                                        actual_x_column = col
                                        break
                            
                            if not actual_x_column:
                                # Fallback: use first column if mapping not available
                                actual_x_column = list(chart_data.columns)[0]
                                print(f"WARNING: Using first column '{actual_x_column}' as X axis (requested '{x_column}' not found)")
                            
                            # Map Y columns - preserve order and all requested columns
                            actual_y_columns = []
                            y_cols_list = y_columns if isinstance(y_columns, list) else [y_columns]
                            
                            for y_col in y_cols_list:
                                actual_y = column_mapping.get(y_col)
                                if not actual_y:
                                    # Try direct match in DataFrame
                                    y_col_normalized = str(y_col).strip().lower()
                                    for col in chart_data.columns:
                                        if str(col).strip().lower() == y_col_normalized and col != actual_x_column:
                                            actual_y = col
                                            break
                                
                                if actual_y and actual_y != actual_x_column:
                                    if actual_y not in actual_y_columns:  # Avoid duplicates
                                        actual_y_columns.append(actual_y)
                            
                            # If no Y columns found, use all columns except X
                            if not actual_y_columns:
                                actual_y_columns = [col for col in chart_data.columns if col != actual_x_column]
                                print(f"WARNING: No Y columns matched, using all columns except X: {actual_y_columns}")
                            
                            print(f"DEBUG: Chart - Requested x_column: '{x_column}', Using: '{actual_x_column}'")
                            print(f"DEBUG: Chart - Requested y_columns: {y_columns}, Using: {actual_y_columns}")
                            print(f"DEBUG: Chart - Data has {len(chart_data)} rows")
                            
                            if actual_x_column and actual_y_columns and len(chart_data) > 0:
                                try:
                                    self.builder.add_chart(
                                        slide, chart_data,
                                        chart_type=chart_config.get("type", "column"),
                                        left=1, top=top_offset, width=8, height=4,
                                        x_column=actual_x_column,  # Use actual column name from DataFrame
                                        y_columns=actual_y_columns,  # Use actual column names from DataFrame
                                        title=chart_config.get("title", ""),
                                        formatting=chart_config.get("formatting", {})
                                    )
                                    chart_success = True
                                    print(f"SUCCESS: Chart successfully added to slide {slide_number}")
                                except Exception as e:
                                    import traceback
                                    error_msg = str(e)
                                    print(f"ERROR: Could not add chart: {error_msg}")
                                    print(f"ERROR Traceback: {traceback.format_exc()}")
                                    chart_success = False
                            else:
                                print(f"ERROR: Cannot create chart - X: {actual_x_column}, Y: {actual_y_columns}, Rows: {len(chart_data)}")
                                chart_success = False
                        else:
                            print(f"ERROR: Chart data is empty - rows: {len(chart_data)}, columns: {len(chart_data.columns)}")
                            chart_success = False
                    else:
                        print(f"ERROR: Chart data retrieval returned None. data_source: {chart_data_source}, sheet: {chart_sheet}")
                        chart_success = False
                        
                except Exception as e:
                    import traceback
                    error_msg = str(e)
                    print(f"ERROR: Exception while retrieving chart data: {error_msg}")
                    print(f"ERROR Traceback: {traceback.format_exc()}")
                    chart_success = False
            else:
                print(f"ERROR: Chart enabled but missing required columns. x_column: {x_column}, y_columns: {y_columns}")
                chart_success = False
        
        # IMPORTANT: Don't add content if chart is enabled (even if chart failed)
        # User wants: if chart is enabled, show chart only, no content mappings
        if chart_config and chart_config.get("enabled", False):
            if chart_success:
                print(f"DEBUG: Chart successfully added to slide {slide_number}, skipping content mappings")
            else:
                print(f"DEBUG: Chart was enabled but failed, skipping content mappings as requested (chart-only slide)")
            return  # Exit early - chart only slide
        
        # Populate content based on mappings
        if content_mappings:
            self.builder.populate_slide_from_mapping(slide, data, {"shape_mappings": content_mappings})
        elif not title and not subtitle and not (chart_config and chart_config.get("enabled", False)):
            # If no content at all, add a placeholder message
            self.builder.add_text_box(
                slide, "No content configured for this slide",
                left=2, top=3, width=6, height=1,
                formatting={
                    "font_size": 18,
                    "font_color": "#999999",
                    "italic": True,
                    "alignment": "center",
                    "font_name": "Calibri"
                }
            )
    
    def _generate_table_slide(self, slide, slide_config: Dict, data: Dict[str, Any], slide_number: int = 1):
        """Generate a table slide."""
        title = slide_config.get("title", "")
        subtitle = slide_config.get("subtitle", "")
        table_mapping = slide_config.get("table_mapping", {})
        chart_config = slide_config.get("chart", None)
        
        # Get formatting from config
        title_formatting = slide_config.get("title_formatting", {})
        subtitle_formatting = slide_config.get("subtitle_formatting", {})
        
        # Default title formatting - use defaults only if user hasn't specified
        default_title_formatting = {
            "font_size": 14,  # Default: 14
            "bold": True,
            "font_color": "#004E6F",  # Default: #004E6F
            "font_name": "Calibri",  # Default: Calibri
            "alignment": "left"  # Left align
        }
        
        # Default subtitle formatting - use defaults only if user hasn't specified
        default_subtitle_formatting = {
            "font_size": 18,  # Default: 18
            "bold": False,
            "font_color": "#009CDE",  # Default: #009CDE
            "font_name": "Georgia",  # Default: Georgia
            "alignment": "left"  # Left align
        }
        
        # Merge: user's formatting overrides defaults
        title_formatting = {**default_title_formatting, **title_formatting}
        subtitle_formatting = {**default_subtitle_formatting, **subtitle_formatting}
        
        # Clean title and subtitle text
        title = self._clean_text(title)
        subtitle = self._clean_text(subtitle)
        
        # Try to set title in existing placeholder
        title_found = False
        if title:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_name = shape.name.lower()
                    if "title" in shape_name or (shape_name == "" and not title_found):
                        shape.text_frame.text = title
                        # Enable word wrap for title
                        shape.text_frame.word_wrap = True
                        # Format title with colors, font, and alignment
                        for paragraph in shape.text_frame.paragraphs:
                            # Set alignment
                            if title_formatting.get("alignment") == "left":
                                paragraph.alignment = PP_ALIGN.LEFT
                            for run in paragraph.runs:
                                # Set font size
                                run.font.size = Pt(title_formatting.get("font_size", 14))
                                # Set font name
                                if "font_name" in title_formatting:
                                    run.font.name = title_formatting["font_name"]
                                # Set bold
                                run.font.bold = title_formatting.get("bold", True)
                                # Apply color
                                if "font_color" in title_formatting:
                                    color_str = title_formatting["font_color"]
                                    if color_str.startswith("#"):
                                        color_str = color_str[1:]
                                    r, g, b = int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16)
                                    from pptx.dml.color import RGBColor
                                    run.font.color.rgb = RGBColor(r, g, b)
                        title_found = True
                        break
        
        # Try to set subtitle in existing placeholder
        subtitle_found = False
        if subtitle:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_name = shape.name.lower()
                    if "subtitle" in shape_name and not subtitle_found:
                        shape.text_frame.text = subtitle
                        # Enable word wrap for subtitle
                        shape.text_frame.word_wrap = True
                        # Format subtitle with colors, font, and alignment
                        for paragraph in shape.text_frame.paragraphs:
                            # Set alignment
                            if subtitle_formatting.get("alignment") == "left":
                                paragraph.alignment = PP_ALIGN.LEFT
                            for run in paragraph.runs:
                                # Set font size
                                run.font.size = Pt(subtitle_formatting.get("font_size", 18))
                                # Set font name
                                if "font_name" in subtitle_formatting:
                                    run.font.name = subtitle_formatting["font_name"]
                                # Set bold
                                run.font.bold = subtitle_formatting.get("bold", False)
                                # Apply color
                                if "font_color" in subtitle_formatting:
                                    color_str = subtitle_formatting["font_color"]
                                    if color_str.startswith("#"):
                                        color_str = color_str[1:]
                                    r, g, b = int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16)
                                    from pptx.dml.color import RGBColor
                                    run.font.color.rgb = RGBColor(r, g, b)
                        subtitle_found = True
                        break
        
        # Calculate table position based on whether we have title/subtitle/chart
        top_offset = 0.5
        if title:
            top_offset = 1.5
        if subtitle:
            top_offset = 2.2
        
        # If no placeholder found, add title as text box
        if not title_found and title:
            # Ensure alignment is set for title
            title_formatting["alignment"] = "left"
            self.builder.add_text_box(
                slide, title,
                left=0.5, top=0.5, width=9, height=0.8,
                formatting=title_formatting
            )
            top_offset = 1.5
        
        # Add subtitle if present and not found in placeholder
        if subtitle and not subtitle_found:
            # Ensure alignment is set for subtitle
            subtitle_formatting["alignment"] = "left"
            self.builder.add_text_box(
                slide, subtitle,
                left=0.5, top=1.5, width=9, height=0.5,
                formatting=subtitle_formatting
            )
            top_offset = 2.2
        
        # Track if chart was successfully added (initialize to False)
        chart_success = False
        
        # Add chart if configured (before table)
        if chart_config and chart_config.get("enabled", False):
            print(f"DEBUG: Chart is enabled for slide {slide_number}")
            # Build proper chart mapping with all necessary columns (x + y)
            x_column = chart_config.get("x_column")
            y_columns = chart_config.get("y_columns", [])
            
            if x_column and y_columns:
                print(f"DEBUG: Chart config - x_column: '{x_column}', y_columns: {y_columns}")
                
                # Use same data source as table if not specified in chart config
                chart_data_source = chart_config.get("data_source") or table_mapping.get("data_source")
                chart_sheet = chart_config.get("sheet") or table_mapping.get("sheet")
                chart_header_row = chart_config.get("header_row", table_mapping.get("header_row", 0))
                
                # Create a mapping for chart data that includes all needed columns
                chart_mapping = {
                    "data_source": chart_data_source,
                    "sheet": chart_sheet,
                    "header_row": chart_header_row,
                    "columns": [x_column] + (y_columns if isinstance(y_columns, list) else [y_columns])
                }
                
                print(f"DEBUG: Chart mapping - data_source: {chart_mapping['data_source']}, sheet: {chart_mapping['sheet']}")
                print(f"DEBUG: Chart mapping - columns: {chart_mapping['columns']}")
                
                try:
                    chart_data = self.builder._get_table_data(data, chart_mapping, return_column_mapping=True)
                    
                    if chart_data is not None:
                        if isinstance(chart_data, tuple):
                            # Returned (DataFrame, column_mapping)
                            chart_data_df, column_mapping = chart_data
                            chart_data = chart_data_df
                        else:
                            # Just DataFrame, create mapping from column order
                            column_mapping = {}
                            actual_columns = list(chart_data.columns)
                            requested_columns = [x_column] + (y_columns if isinstance(y_columns, list) else [y_columns])
                            for i, req_col in enumerate(requested_columns):
                                if i < len(actual_columns):
                                    column_mapping[req_col] = actual_columns[i]
                        
                        print(f"DEBUG: Chart data retrieved - shape: {chart_data.shape}, columns: {list(chart_data.columns)}")
                        print(f"DEBUG: Column mapping: {column_mapping}")
                        
                        if len(chart_data) > 0 and len(chart_data.columns) > 0:
                            # Use column mapping to find actual X and Y column names
                            actual_x_column = column_mapping.get(x_column)
                            if not actual_x_column:
                                # Try to find it in the DataFrame directly
                                x_col_normalized = str(x_column).strip().lower()
                                for col in chart_data.columns:
                                    if str(col).strip().lower() == x_col_normalized:
                                        actual_x_column = col
                                        break
                            
                            if not actual_x_column:
                                # Fallback: use first column if mapping not available
                                actual_x_column = list(chart_data.columns)[0]
                                print(f"WARNING: Using first column '{actual_x_column}' as X axis (requested '{x_column}' not found)")
                            
                            # Map Y columns - preserve order and all requested columns
                            actual_y_columns = []
                            y_cols_list = y_columns if isinstance(y_columns, list) else [y_columns]
                            
                            for y_col in y_cols_list:
                                actual_y = column_mapping.get(y_col)
                                if not actual_y:
                                    # Try direct match in DataFrame
                                    y_col_normalized = str(y_col).strip().lower()
                                    for col in chart_data.columns:
                                        if str(col).strip().lower() == y_col_normalized and col != actual_x_column:
                                            actual_y = col
                                            break
                                
                                if actual_y and actual_y != actual_x_column:
                                    if actual_y not in actual_y_columns:  # Avoid duplicates
                                        actual_y_columns.append(actual_y)
                            
                            # If no Y columns found, use all columns except X
                            if not actual_y_columns:
                                actual_y_columns = [col for col in chart_data.columns if col != actual_x_column]
                                print(f"WARNING: No Y columns matched, using all columns except X: {actual_y_columns}")
                            
                            print(f"DEBUG: Chart - Requested x_column: '{x_column}', Using: '{actual_x_column}'")
                            print(f"DEBUG: Chart - Requested y_columns: {y_columns}, Using: {actual_y_columns}")
                            print(f"DEBUG: Chart - Data has {len(chart_data)} rows")
                            
                            if actual_x_column and actual_y_columns and len(chart_data) > 0:
                                try:
                                    self.builder.add_chart(
                                        slide, chart_data,
                                        chart_type=chart_config.get("type", "column"),
                                        left=1, top=top_offset, width=8, height=3.5,
                                        x_column=actual_x_column,  # Use actual column name from DataFrame
                                        y_columns=actual_y_columns,  # Use actual column names from DataFrame
                                        title=chart_config.get("title", ""),
                                        formatting=chart_config.get("formatting", {})
                                    )
                                    chart_success = True
                                    print(f"SUCCESS: Chart successfully added to slide {slide_number}")
                                except Exception as e:
                                    import traceback
                                    error_msg = str(e)
                                    print(f"ERROR: Could not add chart: {error_msg}")
                                    print(f"ERROR Traceback: {traceback.format_exc()}")
                                    chart_success = False
                            else:
                                print(f"ERROR: Cannot create chart - X: {actual_x_column}, Y: {actual_y_columns}, Rows: {len(chart_data)}")
                                chart_success = False
                        else:
                            print(f"ERROR: Chart data is empty - rows: {len(chart_data)}, columns: {len(chart_data.columns)}")
                            chart_success = False
                    else:
                        print(f"ERROR: Chart data retrieval returned None. data_source: {chart_data_source}, sheet: {chart_sheet}")
                        chart_success = False
                        
                except Exception as e:
                    import traceback
                    error_msg = str(e)
                    print(f"ERROR: Exception while retrieving chart data: {error_msg}")
                    print(f"ERROR Traceback: {traceback.format_exc()}")
                    chart_success = False
            else:
                print(f"ERROR: Chart enabled but missing required columns. x_column: {x_column}, y_columns: {y_columns}")
                chart_success = False
        
        # IMPORTANT: Don't add table if chart is enabled (even if chart failed)
        # User wants: if chart is enabled, show chart only, no table
        if chart_config and chart_config.get("enabled", False):
            if chart_success:
                print(f"DEBUG: Chart successfully added to slide {slide_number}, skipping table generation")
            else:
                print(f"DEBUG: Chart was enabled but failed, skipping table generation as requested (chart-only slide)")
                # Show professional message instead of error
                try:
                    self.builder.add_text_box(
                        slide, "Chart data unavailable",
                        left=2, top=top_offset + 1, width=6, height=1,
                        formatting={
                            "font_size": 14,
                            "font_color": "#666666",
                            "alignment": "center",
                            "font_name": "Calibri"
                        }
                    )
                except:
                    pass  # If message fails, continue
            return  # Exit early, don't add table - chart only slide
        
        # Get table data - only if chart is not enabled or chart failed
        table_data = None
        if table_mapping and table_mapping.get("data_source"):
            table_data = self.builder._get_table_data(data, table_mapping, return_column_mapping=False)
            
            # Handle tuple return if it happens
            if isinstance(table_data, tuple):
                table_data, _ = table_data
        
        # Debug output
        if table_data is not None:
            print(f"DEBUG: Table data retrieved: {len(table_data)} rows, columns: {list(table_data.columns)}")
        else:
            print(f"DEBUG: No table data retrieved for mapping: {table_mapping.get('data_source', 'N/A')} / {table_mapping.get('sheet', 'N/A')}")
        
        if table_data is not None and len(table_data) > 0:
            # Find existing table or add new one
            table_shape = None
            for shape in slide.shapes:
                if shape.has_table:
                    table_shape = shape
                    break
            
            # Log table data information
            num_rows = len(table_data)
            num_cols = len(table_data.columns)
            print(f"DEBUG: Table data prepared - {num_rows} rows, {num_cols} columns")
            print(f"DEBUG: Table columns (in order): {list(table_data.columns)}")
            
            if table_shape:
                self.builder.update_table_data(slide, slide.shapes.index(table_shape), table_data, 
                                             table_mapping.get("formatting"))
            else:
                # Calculate table dimensions - ensure all columns are visible and fit within slide
                # Slide dimensions: 10 inches wide x 7.5 inches tall (standard PowerPoint)
                slide_width = 10.0
                slide_height = 7.5
                
                # Available width with margins
                left_margin = 0.5
                right_margin = 0.5
                available_width = slide_width - left_margin - right_margin  # 9.0 inches
                
                # Available height - account for title/subtitle space
                top_margin = top_offset  # Space for title/subtitle
                bottom_margin = 0.5
                available_height = slide_height - top_margin - bottom_margin
                
                # Calculate optimal width per column - ensure it fits
                if num_cols > 0:
                    # Calculate base width per column
                    base_col_width = available_width / num_cols
                    
                    # Adaptive column width based on number of columns
                    if num_cols <= 4:
                        min_col_width = 1.2  # Wider columns for fewer columns
                        max_col_width = 2.5
                    elif num_cols <= 7:
                        min_col_width = 0.9
                        max_col_width = 2.0
                    elif num_cols <= 10:
                        min_col_width = 0.7
                        max_col_width = 1.5
                    else:
                        min_col_width = 0.5  # Very narrow for many columns
                        max_col_width = 1.2
                    
                    # Ensure columns fit within available width
                    optimal_col_width = min(max(base_col_width, min_col_width), max_col_width)
                    table_width = optimal_col_width * num_cols
                    
                    # Force table to fit within available width (all columns must be visible)
                    if table_width > available_width:
                        optimal_col_width = available_width / num_cols
                        table_width = available_width
                        print(f"INFO: Adjusted column width to {optimal_col_width:.2f} inches to fit {num_cols} columns within slide")
                else:
                    table_width = available_width
                
                # Calculate table height based on rows
                # Use smaller row heights for better fit
                row_height = 0.35  # Reduced row height for more rows
                header_height = 0.5  # Header row slightly taller
                
                # Calculate total height needed
                calculated_height = header_height + (num_rows * row_height)
                
                # Cap table height to available space
                max_table_height = min(available_height, 6.0)  # Don't exceed available height or 6 inches
                table_height = min(max_table_height, max(1.5, calculated_height))
                
                # Position table - ensure it fits
                table_left = left_margin
                
                # Ensure table doesn't go below slide
                max_top = slide_height - table_height - bottom_margin
                table_top = min(top_margin, max_top)
                
                # Ensure minimum spacing between title/subtitle and table
                if table_top < top_offset + 0.2:
                    table_top = top_offset + 0.2
                
                # Final safety check - ensure table fits
                if table_top + table_height > slide_height - bottom_margin:
                    table_height = slide_height - table_top - bottom_margin
                    print(f"INFO: Adjusted table height to {table_height:.2f} inches to fit within slide")
                
                print(f"DEBUG: Table positioning - left: {table_left:.2f}, top: {table_top:.2f}, width: {table_width:.2f}, height: {table_height:.2f}")
                print(f"DEBUG: Table fits within slide bounds: {table_left + table_width <= slide_width}, {table_top + table_height <= slide_height}")
                
                print(f"DEBUG: Table dimensions - width: {table_width:.2f}, height: {table_height:.2f}, columns: {num_cols}")
                print(f"DEBUG: Table position - left: {table_left}, top: {table_top}")
                
                # Add new table with better positioning - ensure all columns are included
                self.builder.add_table(
                    slide, table_data, 
                    left=table_left, top=table_top, 
                    width=table_width, height=table_height,
                    formatting=table_mapping.get("formatting")
                )
                
                print(f"SUCCESS: Table added with {num_cols} columns and {num_rows} rows")
        else:
            # No data available - add informative message with more details
            data_source = table_mapping.get("data_source", "N/A")
            sheet_name = table_mapping.get("sheet", "N/A")
            columns = table_mapping.get("columns", [])
            
            print(f"WARNING: No data available for table. data_source: {data_source}, sheet: {sheet_name}")
            
            # Create empty DataFrame with proper structure for professional display
            if columns and len(columns) > 0:
                # Use requested columns for structure
                empty_data = pd.DataFrame({col: ["No data available"] for col in columns})
            else:
                # Use generic structure
                empty_data = pd.DataFrame({"Data": ["No data available"]})
            
            # Calculate table dimensions for empty state
            num_rows = 1  # One row with message
            num_cols = len(empty_data.columns)
            
            min_col_width = 1.0
            max_col_width = 2.8
            available_width = 9.0
            
            if num_cols > 0:
                optimal_col_width = min(max(available_width / num_cols, min_col_width), max_col_width)
                table_width = optimal_col_width * num_cols
                table_width = min(table_width, available_width)
            else:
                table_width = available_width
            
            row_height = 0.45
            header_height = 0.5
            table_height = header_height + (num_rows * row_height)
            table_height = min(5.0, max(2.5, table_height))
            
            table_left = 0.5
            max_top = 7.0 - table_height
            table_top = min(top_offset, max_top)
            
            if table_top < top_offset + 0.2:
                table_top = top_offset + 0.2
            
            # Add table with empty state - this will show headers and "No data available" message
            self.builder.add_table(
                slide, empty_data,
                left=table_left, top=table_top,
                width=table_width, height=table_height,
                formatting=table_mapping.get("formatting")
            )
    
    def _generate_bullet_list_slide(self, slide, slide_config: Dict, data: Dict[str, Any]):
        """Generate a bullet list slide."""
        title = slide_config.get("title", "")
        items = slide_config.get("items", [])
        
        # Get items from data if specified
        if "items_data_source" in slide_config:
            items = self._get_list_from_data(data, slide_config["items_data_source"])
        
        # Set title
        if title:
            for shape in slide.shapes:
                if shape.has_text_frame and ("Title" in shape.name or shape.name == ""):
                    shape.text_frame.text = title
                    break
        
        # Add bullet list
        self.builder.add_bullet_list(slide, items, 1, 2, 8, 5, 
                                    slide_config.get("formatting"))
    
    def _generate_generic_slide(self, slide, slide_config: Dict, data: Dict[str, Any]):
        """Generate a generic slide based on mappings."""
        mappings = slide_config.get("mappings", {})
        self.builder.populate_slide_from_mapping(slide, data, mappings)
    
    def _get_text_from_data(self, data: Dict[str, Any], data_source_config: Dict) -> str:
        """Extract text value from data."""
        data_source = data_source_config.get("source")
        column = data_source_config.get("column")
        aggregate = data_source_config.get("aggregate", "first")
        
        if data_source in data:
            df = data[data_source]
            if isinstance(df, dict):
                df = list(df.values())[0]
            
            if isinstance(df, pd.DataFrame) and column in df.columns:
                if aggregate == "sum":
                    return str(df[column].sum())
                elif aggregate == "mean":
                    return str(df[column].mean())
                else:
                    return str(df[column].iloc[0])
        
        return data_source_config.get("default", "")
    
    def _get_list_from_data(self, data: Dict[str, Any], data_source_config: Dict) -> List[str]:
        """Extract list of items from data."""
        data_source = data_source_config.get("source")
        column = data_source_config.get("column")
        
        if data_source in data:
            df = data[data_source]
            if isinstance(df, dict):
                df = list(df.values())[0]
            
            if isinstance(df, pd.DataFrame) and column in df.columns:
                return df[column].astype(str).tolist()
        
        return data_source_config.get("default", [])


if __name__ == "__main__":
    # Example usage
    import sys
    
    if len(sys.argv) < 3:
        print("Usage: python ppt_generator.py <template_path> <output_path> [slides_config] [formatting_config]")
        sys.exit(1)
    
    template_path = sys.argv[1]
    output_path = sys.argv[2]
    slides_config = sys.argv[3] if len(sys.argv) > 3 else None
    formatting_config = sys.argv[4] if len(sys.argv) > 4 else None
    
    generator = PPTGenerator(template_path, slides_config, formatting_config)
    
    # Create sample data
    import pandas as pd
    sample_data = {
        "main": pd.DataFrame({
            "Category": ["A", "B", "C"],
            "Value": [10, 20, 30]
        })
    }
    
    generator.generate(sample_data, output_path)

