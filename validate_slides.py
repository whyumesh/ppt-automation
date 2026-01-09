"""
Validate Generated Slides vs Manual PPT
Compares Slide 4 (Consent) and Slide 9 (Chronic & Overcalling)
"""

from pptx import Presentation
import pandas as pd

manual_ppt = "Data/Apr 2025/AIL LT - April'25.pptx"
generated_ppt = "output/test_slide9_raw.pptx"

print("=" * 80)
print("VALIDATING GENERATED SLIDES")
print("=" * 80)

# Load presentations
manual_prs = Presentation(manual_ppt)
generated_prs = Presentation(generated_ppt)

print(f"\nManual PPT: {len(manual_prs.slides)} slides")
print(f"Generated PPT: {len(generated_prs.slides)} slides")

def extract_table_data(slide, slide_num):
    """Extract table data from a slide."""
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break
    
    if not table:
        return None
    
    # Extract table data
    data = []
    for i, row in enumerate(table.rows):
        row_data = [cell.text.strip() for cell in row.cells]
        data.append(row_data)
    
    return data

def compare_tables(manual_data, generated_data, slide_name):
    """Compare two tables and report differences."""
    print(f"\n{'=' * 80}")
    print(f"VALIDATING {slide_name}")
    print(f"{'=' * 80}")
    
    if not manual_data:
        print(f"✗ Manual PPT: No table found")
        return False
    
    if not generated_data:
        print(f"✗ Generated PPT: No table found")
        return False
    
    print(f"\nManual PPT Table: {len(manual_data)} rows, {len(manual_data[0]) if manual_data else 0} cols")
    print(f"Generated PPT Table: {len(generated_data)} rows, {len(generated_data[0]) if generated_data else 0} cols")
    
    # Compare header rows
    if len(manual_data) > 0 and len(generated_data) > 0:
        manual_header = manual_data[0]
        generated_header = generated_data[0]
        
        print(f"\nManual Header: {manual_header}")
        print(f"Generated Header: {generated_header}")
        
        # Compare data rows
        min_rows = min(len(manual_data), len(generated_data))
        matches = 0
        mismatches = []
        
        print(f"\n{'=' * 80}")
        print("ROW-BY-ROW COMPARISON")
        print(f"{'=' * 80}")
        
        for i in range(1, min_rows):  # Skip header row
            manual_row = manual_data[i] if i < len(manual_data) else []
            generated_row = generated_data[i] if i < len(generated_data) else []
            
            # Compare first few columns (division name and key metrics)
            compare_cols = min(len(manual_row), len(generated_row), 4)
            
            if compare_cols > 0:
                manual_vals = manual_row[:compare_cols]
                generated_vals = generated_row[:compare_cols]
                
                # Check if they match (allowing for formatting differences)
                match = True
                for j in range(compare_cols):
                    m_val = str(manual_vals[j]).strip()
                    g_val = str(generated_vals[j]).strip()
                    
                    # Try to compare as numbers if possible
                    try:
                        m_num = float(m_val.replace(',', '').replace('%', ''))
                        g_num = float(g_val.replace(',', '').replace('%', '').replace('.0', ''))
                        if abs(m_num - g_num) > 0.01:  # Allow small differences
                            match = False
                    except:
                        # String comparison
                        if m_val.lower() != g_val.lower():
                            match = False
                
                if match:
                    matches += 1
                    print(f"✓ Row {i}: MATCH")
                    print(f"    Manual:   {manual_vals}")
                    print(f"    Generated: {generated_vals}")
                else:
                    mismatches.append(i)
                    print(f"✗ Row {i}: MISMATCH")
                    print(f"    Manual:   {manual_vals}")
                    print(f"    Generated: {generated_vals}")
        
        print(f"\n{'=' * 80}")
        print("VALIDATION SUMMARY")
        print(f"{'=' * 80}")
        print(f"Total rows compared: {min_rows - 1}")
        print(f"Matches: {matches}")
        print(f"Mismatches: {len(mismatches)}")
        
        if len(mismatches) > 0:
            print(f"\nMismatched rows: {mismatches}")
        
        accuracy = (matches / (min_rows - 1) * 100) if (min_rows - 1) > 0 else 0
        print(f"\nAccuracy: {accuracy:.1f}%")
        
        if accuracy >= 80:
            print(f"✓ VALIDATION PASSED: {slide_name} is mostly correct!")
        elif accuracy >= 50:
            print(f"⚠ VALIDATION PARTIAL: {slide_name} has some issues")
        else:
            print(f"✗ VALIDATION FAILED: {slide_name} needs fixes")
        
        return accuracy >= 80

# Validate Slide 4 (Consent)
if len(manual_prs.slides) > 3 and len(generated_prs.slides) > 3:
    manual_slide4 = manual_prs.slides[3]
    generated_slide4 = generated_prs.slides[3]
    
    manual_data4 = extract_table_data(manual_slide4, 4)
    generated_data4 = extract_table_data(generated_slide4, 4)
    
    slide4_valid = compare_tables(manual_data4, generated_data4, "SLIDE 4 (CONSENT)")
else:
    print("\n✗ Slide 4 not found in one or both presentations")
    slide4_valid = False

# Validate Slide 9 (Chronic & Overcalling)
if len(manual_prs.slides) > 8 and len(generated_prs.slides) > 8:
    manual_slide9 = manual_prs.slides[8]
    generated_slide9 = generated_prs.slides[8]
    
    manual_data9 = extract_table_data(manual_slide9, 9)
    generated_data9 = extract_table_data(generated_slide9, 9)
    
    slide9_valid = compare_tables(manual_data9, generated_data9, "SLIDE 9 (CHRONIC & OVERCALLING)")
else:
    print("\n✗ Slide 9 not found in one or both presentations")
    slide9_valid = False

# Overall summary
print(f"\n{'=' * 80}")
print("OVERALL VALIDATION SUMMARY")
print(f"{'=' * 80}")
print(f"Slide 4 (Consent): {'✓ PASS' if slide4_valid else '✗ FAIL'}")
print(f"Slide 9 (Chronic & Overcalling): {'✓ PASS' if slide9_valid else '✗ FAIL'}")

if slide4_valid and slide9_valid:
    print("\n🎉 ALL VALIDATIONS PASSED!")
elif slide4_valid or slide9_valid:
    print("\n⚠ PARTIAL SUCCESS - Some slides need fixes")
else:
    print("\n✗ VALIDATION FAILED - Slides need fixes")

